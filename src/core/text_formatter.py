# -*- coding: utf-8 -*-
"""
文本格式化器
负责处理文本中的重点/难点标记
"""

from typing import List, Tuple, Optional
from pptx.dml.color import RGBColor


class TextFormatter:
    """文本格式化器类"""

    # 颜色定义
    HIGHLIGHT_COLOR = "0070C0"  # 重点 - 蓝色
    DIFFICULTY_COLOR = "ED7D31"  # 难点 - 橙色

    def __init__(self):
        """初始化文本格式化器"""
        pass

    def parse_text_with_marks(
        self,
        text: str,
        highlights: List[str],
        difficulties: List[str]
    ) -> List[Tuple[str, Optional[str]]]:
        """
        解析文本并标记重点/难点

        Args:
            text: 原始文本
            highlights: 重点文本列表
            difficulties: 难点文本列表

        Returns:
            列表，每个元素为 (文本片段, 颜色) 元组
            颜色为 None 表示默认颜色
        """
        if not text:
            return []

        # 创建标记列表：(起始位置, 结束位置, 颜色)
        marks = []

        # 收集所有重点标记
        for highlight in highlights:
            if highlight and highlight in text:
                start = 0
                while True:
                    pos = text.find(highlight, start)
                    if pos == -1:
                        break
                    marks.append((pos, pos + len(highlight), self.HIGHLIGHT_COLOR))
                    start = pos + len(highlight)

        # 收集所有难点标记
        for difficulty in difficulties:
            if difficulty and difficulty in text:
                start = 0
                while True:
                    pos = text.find(difficulty, start)
                    if pos == -1:
                        break
                    marks.append((pos, pos + len(difficulty), self.DIFFICULTY_COLOR))
                    start = pos + len(difficulty)

        # 如果没有标记，返回整个文本
        if not marks:
            return [(text, None)]

        # 按起始位置排序
        marks.sort(key=lambda x: x[0])

        # 合并重叠的标记（优先保留第一个标记的颜色）
        merged_marks = []
        for mark in marks:
            if not merged_marks:
                merged_marks.append(mark)
            else:
                last = merged_marks[-1]
                # 检查是否重叠
                if mark[0] < last[1]:
                    # 重叠，保留先出现的标记
                    continue
                else:
                    merged_marks.append(mark)

        # 构建结果列表
        result = []
        last_end = 0

        for start, end, color in merged_marks:
            # 添加前面的普通文本
            if start > last_end:
                result.append((text[last_end:start], None))

            # 添加标记文本
            result.append((text[start:end], color))
            last_end = end

        # 添加剩余的普通文本
        if last_end < len(text):
            result.append((text[last_end:], None))

        return result

    def apply_format_to_run(self, run, color: str):
        """
        将颜色应用到Run对象

        Args:
            run: python-pptx的Run对象
            color: 颜色代码（如 "0070C0"）
        """
        if color:
            run.font.color.rgb = RGBColor.from_string(color)

    def format_text_frame(self, text_frame, text: str,
                          highlights: List[str] = None,
                          difficulties: List[str] = None):
        """
        格式化文本框，应用重点/难点标记

        Args:
            text_frame: python-pptx的TextFrame对象
            text: 要设置的文本
            highlights: 重点文本列表
            difficulties: 难点文本列表
        """
        highlights = highlights or []
        difficulties = difficulties or []

        # 解析文本
        segments = self.parse_text_with_marks(text, highlights, difficulties)

        if not segments:
            return

        # 清除原有内容
        text_frame.clear()

        # 获取第一个段落
        paragraph = text_frame.paragraphs[0]

        # 添加格式化的文本
        first = True
        for segment_text, color in segments:
            if first:
                run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                first = False
            else:
                run = paragraph.add_run()

            run.text = segment_text

            if color:
                self.apply_format_to_run(run, color)

    def format_paragraph(self, paragraph, text: str,
                        highlights: List[str] = None,
                        difficulties: List[str] = None):
        """
        格式化段落，应用重点/难点标记

        Args:
            paragraph: python-pptx的Paragraph对象
            text: 要设置的文本
            highlights: 重点文本列表
            difficulties: 难点文本列表
        """
        highlights = highlights or []
        difficulties = difficulties or []

        # 解析文本
        segments = self.parse_text_with_marks(text, highlights, difficulties)

        if not segments:
            return

        # 清除原有runs
        for run in paragraph.runs:
            run.text = ""

        # 添加格式化的文本
        for i, (segment_text, color) in enumerate(segments):
            if i < len(paragraph.runs):
                run = paragraph.runs[i]
            else:
                run = paragraph.add_run()

            run.text = segment_text

            if color:
                self.apply_format_to_run(run, color)

    @staticmethod
    def get_color_rgb(color_str: str) -> RGBColor:
        """
        将颜色字符串转换为RGBColor对象

        Args:
            color_str: 颜色字符串（如 "0070C0"）

        Returns:
            RGBColor对象
        """
        return RGBColor.from_string(color_str)

    @staticmethod
    def highlight_color() -> RGBColor:
        """获取重点颜色（蓝色）"""
        return RGBColor.from_string(TextFormatter.HIGHLIGHT_COLOR)

    @staticmethod
    def difficulty_color() -> RGBColor:
        """获取难点颜色（橙色）"""
        return RGBColor.from_string(TextFormatter.DIFFICULTY_COLOR)
