# -*- coding: utf-8 -*-
"""
幻灯片内容填充器
负责填充课程信息页的内容
"""

from typing import Dict, Any, Optional, List
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .models import CourseUnitData, EvaluationLevel, OverallEvaluation, HomeworkEvaluation


class SlideContentManager:
    """幻灯片内容管理器"""

    # 复选框字符
    CHECKBOX_UNCHECKED = "□"
    CHECKBOX_CHECKED = "☑"

    # 文本占位符映射（基于模板分析）
    FIELD_MAPPINGS = {
        "lesson_number": ["（第 1 课）", "（第1课）"],
        "course_content": ["机械臂设计（1阶）"],
        "student_name": ["赵如一"],
        "teacher_name": ["于双源"],
        "class_hours": ["2课时"],
        "class_date": ["2025.7.23 18:30-20:30"],
    }

    # 12项评价的字段名（对应models.py中的字段）
    EVALUATION_FIELDS = [
        "logic_thinking",           # 逻辑思维表现
        "content_understanding",    # 课堂内容理解
        "task_completion",          # 课堂任务完成
        "listening_habit",          # 课堂听课习惯
        "problem_solving",          # 解决问题表现
        "independent_analysis",     # 独立分析思考
        "knowledge_proficiency",    # 知识熟练程度
        "imagination_creativity",   # 想象力与创新
        "frustration_handling",     # 挫折困难应对
        "learning_method",          # 学习方法习惯
        "hands_on_ability",         # 动手实践能力
        "focus_efficiency",         # 专注度与效率
    ]

    # 评价等级对应的文本
    EVALUATION_TEXTS = {
        "优": "优",
        "良": "良",
        "中": "中",
        "差": "差",
        "未体现": "未体现"
    }

    def __init__(self, slide):
        """
        初始化内容管理器

        Args:
            slide: pptx Slide对象
        """
        self.slide = slide
        self._shape_map = {}
        self._build_shape_map()

    def _build_shape_map(self):
        """构建形状映射表"""
        self._shape_map = {}

        def process_shape(shape, parent_name=""):
            # 记录有文本的形状
            if hasattr(shape, 'text') and shape.text:
                # 使用文本内容作为键的一部分
                text_key = shape.text.strip()[:50]
                if text_key not in self._shape_map:
                    self._shape_map[text_key] = []
                self._shape_map[text_key].append(shape)

            # 处理组合形状
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for sub_shape in shape.shapes:
                        process_shape(sub_shape, shape.name)
                except:
                    pass

        for shape in self.slide.shapes:
            process_shape(shape)

    def find_shape_by_text(self, text: str):
        """
        根据文本内容查找形状

        Args:
            text: 要查找的文本

        Returns:
            找到的形状，如果没有找到返回None
        """
        text = text.strip()
        for key, shapes in self._shape_map.items():
            if text in key or key in text:
                return shapes[0]
        return None

    def find_shapes_by_keyword(self, keyword: str) -> List:
        """
        根据关键词查找所有匹配的形状

        Args:
            keyword: 关键词

        Returns:
            匹配的形状列表
        """
        result = []
        for key, shapes in self._shape_map.items():
            if keyword in key:
                result.extend(shapes)
        return result

    def replace_text(self, old_text: str, new_text: str) -> bool:
        """
        替换文本

        Args:
            old_text: 旧文本
            new_text: 新文本

        Returns:
            是否替换成功
        """
        shape = self.find_shape_by_text(old_text)
        if shape:
            shape.text = new_text
            return True
        return False

    def replace_text_in_shape(self, shape, old_text: str, new_text: str) -> bool:
        """
        在形状中替换文本（保留格式）

        Args:
            shape: 形状对象
            old_text: 旧文本
            new_text: 新文本

        Returns:
            是否替换成功
        """
        if not hasattr(shape, 'text_frame'):
            return False

        text_frame = shape.text_frame
        full_text = ""

        # 收集所有段落和运行
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                full_text += run.text
            full_text += "\n"

        # 替换文本
        if old_text in full_text:
            new_full_text = full_text.replace(old_text, new_text)

            # 清除原有内容
            text_frame.clear()

            # 添加新内容
            p = text_frame.paragraphs[0]
            p.text = new_full_text.strip()
            return True

        return False

    def fill_basic_info(self, data: CourseUnitData) -> Dict[str, bool]:
        """
        填充基本信息

        Args:
            data: 课程单元数据

        Returns:
            填充结果字典
        """
        results = {}

        # 填充课时编号
        lesson_text = f"（第 {data.lesson_number} 课）"
        results["lesson_number"] = self._fill_field("lesson_number", lesson_text)

        # 填充课程内容
        if data.course_content:
            results["course_content"] = self._fill_field("course_content", data.course_content)

        # 填充学生姓名
        if data.student_name:
            results["student_name"] = self._fill_field("student_name", data.student_name)

        # 填充授课教师
        if data.teacher_name:
            results["teacher_name"] = self._fill_field("teacher_name", data.teacher_name)

        # 填充课时数
        hours_text = f"{data.class_hours}课时"
        results["class_hours"] = self._fill_field("class_hours", hours_text)

        # 填充上课日期
        if data.class_date:
            results["class_date"] = self._fill_field("class_date", data.class_date)

        return results

    def _fill_field(self, field_name: str, new_value: str) -> bool:
        """
        填充指定字段

        Args:
            field_name: 字段名
            new_value: 新值

        Returns:
            是否成功
        """
        if field_name not in self.FIELD_MAPPINGS:
            return False

        for old_value in self.FIELD_MAPPINGS[field_name]:
            shape = self.find_shape_by_text(old_value)
            if shape:
                shape.text = new_value
                return True

        return False

    def fill_course_content(self, content: str, highlights: List[str] = None,
                           difficulties: List[str] = None) -> bool:
        """
        填充课程内容（支持重点和难点标记）

        Args:
            content: 课程内容
            highlights: 重点列表
            difficulties: 难点列表

        Returns:
            是否成功
        """
        # 查找课程内容区域（组合 7 中的矩形）
        shapes = self.find_shapes_by_keyword("理解什么是三维")
        if shapes:
            shapes[0].text = content
            return True

        # 尝试其他方式查找
        for key, shapes in self._shape_map.items():
            if "理解" in key or "借助" in key:
                shapes[0].text = content
                return True

        return False

    def fill_evaluation(self, field_index: int, level: EvaluationLevel) -> bool:
        """
        填充单个评价项

        Args:
            field_index: 评价项索引 (0-11)
            level: 评价等级

        Returns:
            是否成功
        """
        if field_index < 0 or field_index >= 12:
            return False

        # 评价选项文本
        options = ["优", "良", "中", "差", "未体现"]

        # 找到对应的评价行
        # 这里需要根据实际模板结构来实现
        # 查找包含所有选项的形状
        for key, shapes in self._shape_map.items():
            if all(opt in key for opt in options):
                # 这是评价行，根据level设置选中状态
                new_text = self._set_checkbox(key, level.value)
                shapes[0].text = new_text
                return True

        return False

    def _set_checkbox(self, text: str, selected: str) -> str:
        """
        设置复选框选中状态

        Args:
            text: 原始文本
            selected: 选中的选项

        Returns:
            修改后的文本
        """
        options = ["优", "良", "中", "差", "未体现"]

        result = text
        for opt in options:
            if opt == selected:
                # 选中状态
                result = result.replace(f"{opt}{self.CHECKBOX_UNCHECKED}",
                                       f"{opt}{self.CHECKBOX_CHECKED}")
                result = result.replace(f"{opt}", f"{opt}{self.CHECKBOX_CHECKED}")
            else:
                # 未选中状态
                result = result.replace(f"{opt}{self.CHECKBOX_CHECKED}",
                                       f"{opt}{self.CHECKBOX_UNCHECKED}")
                result = result.replace(f"{opt}", f"{opt}{self.CHECKBOX_UNCHECKED}")

        return result

    def fill_all_evaluations(self, data: CourseUnitData) -> Dict[str, bool]:
        """
        填充所有评价项

        Args:
            data: 课程单元数据

        Returns:
            填充结果字典
        """
        results = {}

        # 填充12项评价（字段名与models.py一致）
        eval_fields = [
            ("logic_thinking", data.logic_thinking),
            ("content_understanding", data.content_understanding),
            ("task_completion", data.task_completion),
            ("listening_habit", data.listening_habit),
            ("problem_solving", data.problem_solving),
            ("independent_analysis", data.independent_analysis),
            ("knowledge_proficiency", data.knowledge_proficiency),
            ("imagination_creativity", data.imagination_creativity),
            ("frustration_handling", data.frustration_handling),
            ("learning_method", data.learning_method),
            ("hands_on_ability", data.hands_on_ability),
            ("focus_efficiency", data.focus_efficiency),
        ]

        for i, (field_name, level) in enumerate(eval_fields):
            if level:
                results[field_name] = self.fill_evaluation(i, level)

        return results

    def fill_overall_evaluation(self, level: OverallEvaluation) -> bool:
        """
        填充总体评价

        Args:
            level: 总体评价等级

        Returns:
            是否成功
        """
        options = ["优", "良", "仍需努力", "需要改进"]

        for key, shapes in self._shape_map.items():
            if "总体表现" in key or all(opt in key for opt in ["优", "良", "仍需努力"]):
                new_text = self._set_overall_checkbox(key, level.value, options)
                shapes[0].text = new_text
                return True

        return False

    def _set_overall_checkbox(self, text: str, selected: str, options: List[str]) -> str:
        """设置总体评价复选框"""
        result = text
        for opt in options:
            if opt == selected:
                result = result.replace(f"{opt}", f"{opt}{self.CHECKBOX_CHECKED}")
            else:
                result = result.replace(f"{opt}", f"{opt}{self.CHECKBOX_UNCHECKED}")
        return result

    def fill_homework_evaluation(self, level: HomeworkEvaluation) -> bool:
        """
        填充上次作业评价

        Args:
            level: 作业评价等级

        Returns:
            是否成功
        """
        options = ["优", "良", "中", "差", "无"]

        for key, shapes in self._shape_map.items():
            if "上次作业情况" in key:
                new_text = self._set_homework_checkbox(key, level.value, options)
                shapes[0].text = new_text
                return True

        return False

    def _set_homework_checkbox(self, text: str, selected: str, options: List[str]) -> str:
        """设置作业评价复选框"""
        result = text
        for opt in options:
            if opt == selected:
                result = result.replace(f"{opt}", f"{opt}{self.CHECKBOX_CHECKED}")
            else:
                result = result.replace(f"{opt}", f"{opt}{self.CHECKBOX_UNCHECKED}")
        return result

    def fill_supplementary_notes(self, notes: str) -> bool:
        """
        填充补充说明

        Args:
            notes: 补充说明内容

        Returns:
            是否成功
        """
        # 查找补充说明区域
        for key, shapes in self._shape_map.items():
            if "今天开始" in key or len(key) > 50:
                # 可能是补充说明区域
                shapes[0].text = notes
                return True

        return False

    def fill_homework(self, homework: str) -> bool:
        """
        填充课堂作业

        Args:
            homework: 作业内容

        Returns:
            是否成功
        """
        for key, shapes in self._shape_map.items():
            if "本次课无课堂作业" in key:
                shapes[0].text = homework
                return True

        return False

    def fill_notes(self, notes: str) -> bool:
        """
        填充其它注意事项

        Args:
            notes: 注意事项内容

        Returns:
            是否成功
        """
        for key, shapes in self._shape_map.items():
            if "三维建模对空间想象力" in key or "其它注意事项" in key:
                shapes[0].text = notes
                return True

        return False

    def fill_all(self, data: CourseUnitData) -> Dict[str, bool]:
        """
        填充所有内容

        Args:
            data: 课程单元数据

        Returns:
            填充结果字典
        """
        results = {}

        # 填充基本信息
        results.update(self.fill_basic_info(data))

        # 填充课程内容
        if data.course_content:
            results["course_content_detail"] = self.fill_course_content(
                data.course_content,
                data.highlights,
                data.difficulties
            )

        # 填充评价
        results.update(self.fill_all_evaluations(data))

        # 填充总体评价
        if data.overall_evaluation:
            results["overall_evaluation"] = self.fill_overall_evaluation(data.overall_evaluation)

        # 填充上次作业评价
        if data.last_homework_status:
            results["last_homework_status"] = self.fill_homework_evaluation(data.last_homework_status)

        # 填充补充说明
        if data.additional_comments:
            results["additional_comments"] = self.fill_supplementary_notes(data.additional_comments)

        # 填充课堂作业
        if data.homework:
            results["homework"] = self.fill_homework(data.homework)

        # 填充注意事项
        if data.other_notes:
            results["other_notes"] = self.fill_notes(data.other_notes)

        return results
