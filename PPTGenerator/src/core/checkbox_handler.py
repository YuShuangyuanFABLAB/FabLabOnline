# -*- coding: utf-8 -*-
"""
复选框处理器
负责处理PPT中的评价复选框切换
使用python-pptx原生API，不直接操作XML
"""

from typing import List, Dict, Optional
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor


class CheckboxHandler:
    """复选框处理器类"""

    # 模板中使用的Wingdings字符
    CHECKBOX_CHECKED = '\uf052'   # 选中（实心方块）
    CHECKBOX_UNCHECKED = '\uf0a3'  # 未选中（空心方块）

    # 12项评价的选项
    EVALUATION_OPTIONS = ["优", "良", "中", "差", "未体现"]

    # 总体评价的选项
    OVERALL_OPTIONS = ["优", "良", "仍需努力", "需要改进"]

    # 上次作业评价的选项
    HOMEWORK_OPTIONS = ["优", "良", "中", "差", "无"]

    # 颜色常量
    COLOR_SELECTED = MSO_THEME_COLOR.ACCENT_2  # 橙色
    COLOR_UNSELECTED = MSO_THEME_COLOR.ACCENT_1  # 蓝色

    def __init__(self, slide):
        self.slide = slide
        self._evaluation_shapes = []
        self._overall_shape = None
        self._homework_shape = None
        self._analyze_slide()

    def _analyze_slide(self):
        """分析幻灯片，识别评价相关形状"""
        self._evaluation_shapes = []

        def process_shape(shape):
            if hasattr(shape, 'text') and shape.text:
                text = shape.text
                has_eval = all(opt in text for opt in ["优", "良", "中", "差"])
                has_not_shown = "未体现" in text
                has_none = "无" in text
                has_overall = "仍需努力" in text or "需要改进" in text

                if has_eval and has_not_shown and not has_none and not has_overall:
                    top = shape.top if hasattr(shape, 'top') else 0
                    left = shape.left if hasattr(shape, 'left') else 0
                    self._evaluation_shapes.append((top, left, shape))
                elif has_overall:
                    self._overall_shape = shape
                elif "上次作业情况" in text or (has_eval and has_none and not has_not_shown):
                    self._homework_shape = shape

            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for sub in shape.shapes:
                        process_shape(sub)
                except:
                    pass

        for shape in self.slide.shapes:
            process_shape(shape)

        # 排序：左列从上到下，再右列从上到下
        LEFT_THRESHOLD = 2000000
        left_col = [(t, l, s) for t, l, s in self._evaluation_shapes if l < LEFT_THRESHOLD]
        right_col = [(t, l, s) for t, l, s in self._evaluation_shapes if l >= LEFT_THRESHOLD]
        left_col.sort(key=lambda x: x[0])
        right_col.sort(key=lambda x: x[0])
        self._evaluation_shapes = [(t, s) for t, l, s in left_col + right_col]

    def get_evaluation_count(self) -> int:
        return len(self._evaluation_shapes)

    def set_evaluation(self, index: int, selected_value: str) -> bool:
        if index < 0 or index >= len(self._evaluation_shapes):
            return False
        _, shape = self._evaluation_shapes[index]
        return self._update_checkboxes(shape, selected_value, self.EVALUATION_OPTIONS)

    def set_overall_evaluation(self, selected_value: str) -> bool:
        if not self._overall_shape:
            return False
        return self._update_checkboxes(self._overall_shape, selected_value, self.OVERALL_OPTIONS)

    def set_homework_evaluation(self, selected_value: str) -> bool:
        if not self._homework_shape:
            return False
        return self._update_checkboxes(self._homework_shape, selected_value, self.HOMEWORK_OPTIONS)

    def _update_checkboxes(self, shape, selected: str, options: List[str]) -> bool:
        """更新复选框"""
        if not hasattr(shape, 'text_frame'):
            return False

        from lxml import etree
        NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

        # 判断区域类型
        is_overall = "仍需努力" in options or "需要改进" in options  # 总体评价
        is_homework = "无" in options and "未体现" not in options  # 上次作业情况

        # 第一步：将所有选中框改成未选中框
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text
                if self.CHECKBOX_CHECKED in text:
                    # 替换字符
                    run.text = text.replace(self.CHECKBOX_CHECKED, self.CHECKBOX_UNCHECKED)

                    # 根据区域类型设置未选中框的颜色
                    run_element = run._r
                    rPr = run_element.find(f'{NS}rPr')
                    if rPr is None:
                        rPr = etree.SubElement(run_element, f'{NS}rPr')

                    # 移除旧的颜色设置
                    old_solidFill = rPr.find(f'{NS}solidFill')
                    if old_solidFill is not None:
                        rPr.remove(old_solidFill)

                    # 设置新的颜色（根据区域类型，solidFill必须在sym之前）
                    if is_overall or is_homework:
                        # 查找sym元素的位置
                        sym_elem = rPr.find(f'{NS}sym')

                        # 创建新的solidFill
                        new_solidFill = etree.Element(f'{NS}solidFill')
                        if is_overall:
                            schemeClr = etree.SubElement(new_solidFill, f'{NS}schemeClr')
                            schemeClr.set('val', 'bg1')
                        else:
                            srgbClr = etree.SubElement(new_solidFill, f'{NS}srgbClr')
                            srgbClr.set('val', '36A4A4')

                        # 在sym之前插入solidFill
                        if sym_elem is not None:
                            sym_index = list(rPr).index(sym_elem)
                            rPr.insert(sym_index, new_solidFill)
                        else:
                            rPr.append(new_solidFill)
                    # else: 课堂表现 - 不设置颜色，保持继承

        # 第二步：将选中选项的未选中框改成选中框，并设置选中颜色
        for para in shape.text_frame.paragraphs:
            self._select_and_color_checkbox(para, selected, options)

        return True

    def _select_and_color_checkbox(self, paragraph, selected: str, options: List[str]):
        """在段落中选中指定的选项，并设置选中框颜色"""
        from lxml import etree

        # 判断是否是总体评价区域：选项包含"仍需努力"或"需要改进"
        is_overall = "仍需努力" in options or "需要改进" in options

        # 收集所有run及其文本
        runs_info = []
        for run in paragraph.runs:
            runs_info.append({
                'run': run,
                'text': run.text,
                'element': run._r
            })

        # 合并所有文本
        full_text = ''.join(info['text'] for info in runs_info)

        # 找到需要选中的位置
        pattern = selected + self.CHECKBOX_UNCHECKED
        pos = full_text.find(pattern)
        if pos == -1:
            return

        # 计算复选框字符在合并文本中的绝对位置
        checkbox_pos = pos + len(selected)

        # 找到包含该位置的run
        current_pos = 0
        for info in runs_info:
            run_text = info['text']
            run_start = current_pos
            run_end = current_pos + len(run_text)

            if run_start <= checkbox_pos < run_end:
                # 复选框在这个run中
                local_pos = checkbox_pos - run_start
                if run_text[local_pos] == self.CHECKBOX_UNCHECKED:
                    # 获取原始run的rPr
                    run_element = info['run']._r
                    rPr_element = run_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')

                    # 复制原始rPr的属性
                    original_attrs = {}
                    original_sym = None
                    if rPr_element is not None:
                        for attr in rPr_element.attrib:
                            original_attrs[attr] = rPr_element.get(attr)
                        sym = rPr_element.find('{http://schemas.openxmlformats.org/drawingml/2006/main}sym')
                        if sym is not None:
                            original_sym = dict(sym.attrib)

                    # 检查run是否只有这一个方框
                    if run_text == self.CHECKBOX_UNCHECKED:
                        # 直接修改为选中框，并设置颜色
                        info['run'].text = self.CHECKBOX_CHECKED
                        # 根据区域设置不同颜色（使用RGB颜色，避免theme_color的问题）
                        NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                        rPr = info['run']._r.get_or_add_rPr()

                        # 移除现有的颜色设置
                        for color_elem in rPr.findall(f'{NS}solidFill') + rPr.findall(f'{NS}schemeClr'):
                            rPr.remove(color_elem)

                        # 查找sym元素的位置，确保solidFill在sym之前
                        sym_elem = rPr.find(f'{NS}sym')

                        # 创建新的颜色元素
                        solidFill = etree.Element(f'{NS}solidFill')
                        srgbClr = etree.SubElement(solidFill, f'{NS}srgbClr')
                        if is_overall:
                            srgbClr.set('val', 'FFFF00')  # 黄色
                        else:
                            srgbClr.set('val', 'ED7D31')  # 橙色

                        # 在sym之前插入solidFill（如果sym存在）
                        if sym_elem is not None:
                            sym_index = list(rPr).index(sym_elem)
                            rPr.insert(sym_index, solidFill)
                        else:
                            rPr.append(solidFill)
                    else:
                        # 需要拆分run
                        before = run_text[:local_pos]
                        checkbox = self.CHECKBOX_CHECKED
                        after = run_text[local_pos + 1:]

                        NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

                        # 处理原始run（before部分）
                        # 如果before不包含复选框字符，需要移除sym元素
                        has_checkbox_in_before = self.CHECKBOX_UNCHECKED in before or self.CHECKBOX_CHECKED in before
                        if before:
                            info['run'].text = before
                            if not has_checkbox_in_before:
                                # 移除sym元素（普通文字不需要sym）
                                run_rPr = run_element.find(f'{NS}rPr')
                                if run_rPr is not None:
                                    sym_elem = run_rPr.find(f'{NS}sym')
                                    if sym_elem is not None:
                                        run_rPr.remove(sym_elem)
                        else:
                            info['run'].text = ""

                        # 创建选中框run
                        p_element = paragraph._p
                        run_index = list(p_element).index(run_element)

                        checked_r = etree.Element(f'{NS}r')
                        checked_rPr = etree.SubElement(checked_r, f'{NS}rPr')

                        # 复制原始属性
                        for attr, val in original_attrs.items():
                            checked_rPr.set(attr, val)

                        # 先设置颜色（solidFill必须在sym之前）
                        solidFill = etree.SubElement(checked_rPr, f'{NS}solidFill')
                        srgbClr = etree.SubElement(solidFill, f'{NS}srgbClr')
                        if is_overall:
                            srgbClr.set('val', 'FFFF00')  # 黄色（总体评价）
                        else:
                            srgbClr.set('val', 'ED7D31')  # 橙色（课堂表现/作业评价）

                        # 然后添加sym元素
                        if original_sym:
                            checked_sym = etree.SubElement(checked_rPr, f'{NS}sym')
                            for attr, val in original_sym.items():
                                checked_sym.set(attr, val)

                        # 设置文本
                        checked_t = etree.SubElement(checked_r, f'{NS}t')
                        checked_t.text = checkbox

                        # 插入选中框run
                        run_index += 1
                        p_element.insert(run_index, checked_r)

                        # 如果有后部，创建后部run
                        if after:
                            after_r = etree.Element(f'{NS}r')
                            after_rPr = etree.SubElement(after_r, f'{NS}rPr')

                            for attr, val in original_attrs.items():
                                after_rPr.set(attr, val)

                            # 先复制原始颜色（solidFill必须在sym之前）
                            if rPr_element is not None:
                                orig_solidFill = rPr_element.find(f'{NS}solidFill')
                                if orig_solidFill is not None:
                                    new_solidFill = etree.SubElement(after_rPr, f'{NS}solidFill')
                                    for child in orig_solidFill:
                                        if child.tag == f'{NS}schemeClr':
                                            schemeClr = etree.SubElement(new_solidFill, f'{NS}schemeClr')
                                            for attr in child.attrib:
                                                schemeClr.set(attr, child.get(attr))
                                        elif child.tag == f'{NS}srgbClr':
                                            srgbClr = etree.SubElement(new_solidFill, f'{NS}srgbClr')
                                            for attr in child.attrib:
                                                srgbClr.set(attr, child.get(attr))

                            # 然后添加sym元素（只有当after包含复选框字符时）
                            has_checkbox_in_after = self.CHECKBOX_UNCHECKED in after or self.CHECKBOX_CHECKED in after
                            if original_sym and has_checkbox_in_after:
                                after_sym = etree.SubElement(after_rPr, f'{NS}sym')
                                for attr, val in original_sym.items():
                                    after_sym.set(attr, val)
                                orig_solidFill = rPr_element.find(f'{NS}solidFill')
                                if orig_solidFill is not None:
                                    new_solidFill = etree.SubElement(after_rPr, f'{NS}solidFill')
                                    for child in orig_solidFill:
                                        if child.tag == f'{NS}schemeClr':
                                            schemeClr = etree.SubElement(new_solidFill, f'{NS}schemeClr')
                                            for attr in child.attrib:
                                                schemeClr.set(attr, child.get(attr))
                                        elif child.tag == f'{NS}srgbClr':
                                            srgbClr = etree.SubElement(new_solidFill, f'{NS}srgbClr')
                                            for attr in child.attrib:
                                                srgbClr.set(attr, child.get(attr))

                            after_t = etree.SubElement(after_r, f'{NS}t')
                            after_t.text = after
                            if after and (after[0] == ' ' or after[-1] == ' '):
                                after_t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')

                            run_index += 1
                            p_element.insert(run_index, after_r)
                return

            current_pos = run_end

    def _select_option_in_paragraph(self, paragraph, selected: str, options: List[str]):
        """在段落中选中指定的选项"""
        from lxml import etree

        # 收集所有run及其文本
        runs_info = []
        for run in paragraph.runs:
            runs_info.append({
                'run': run,
                'text': run.text,
                'element': run._r
            })

        # 合并所有文本
        full_text = ''.join(info['text'] for info in runs_info)

        # 找到需要选中的位置
        pattern = selected + self.CHECKBOX_UNCHECKED
        pos = full_text.find(pattern)
        if pos == -1:
            return

        # 计算复选框字符在合并文本中的绝对位置
        checkbox_pos = pos + len(selected)

        # 找到包含该位置的run
        current_pos = 0
        for info in runs_info:
            run_text = info['text']
            run_start = current_pos
            run_end = current_pos + len(run_text)

            if run_start <= checkbox_pos < run_end:
                # 复选框在这个run中
                local_pos = checkbox_pos - run_start
                if run_text[local_pos] == self.CHECKBOX_UNCHECKED:
                    # 替换字符
                    new_text = run_text[:local_pos] + self.CHECKBOX_CHECKED + run_text[local_pos + 1:]
                    info['run'].text = new_text
                return

            current_pos = run_end

    def _apply_checkbox_colors(self, shape, selected: str, options: List[str]):
        """为复选框应用颜色"""
        if not hasattr(shape, 'text_frame'):
            return

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = run.text
                if not text:
                    continue

                # 检查这个run是否包含复选框
                has_checked = self.CHECKBOX_CHECKED in text
                has_unchecked = self.CHECKBOX_UNCHECKED in text

                if has_checked or has_unchecked:
                    # 如果run只包含复选框，直接设置颜色
                    if text in [self.CHECKBOX_CHECKED, self.CHECKBOX_UNCHECKED]:
                        if text == self.CHECKBOX_CHECKED:
                            run.font.color.theme_color = self.COLOR_SELECTED
                        else:
                            run.font.color.theme_color = self.COLOR_UNSELECTED
                    else:
                        # run包含多个字符，需要拆分
                        self._split_and_color_run(para, run, selected)

    def _split_and_color_run(self, paragraph, run, selected: str):
        """拆分包含复选框的run并设置颜色"""
        from lxml import etree

        text = run.text
        if not text:
            return

        # 获取run的原始格式信息（从XML中读取）
        run_element = run._r
        rPr_element = run_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')

        # 检查原始run是否使用sym字体（符号字体如Wingdings 2）
        original_sym = None
        original_sz = None
        original_lang = None

        if rPr_element is not None:
            # 查找sym元素
            original_sym = rPr_element.find('{http://schemas.openxmlformats.org/drawingml/2006/main}sym')
            # 获取字号
            original_sz = rPr_element.get('sz')
            # 获取语言
            original_lang = rPr_element.get('lang')

        # 找出所有复选框的位置
        segments = []  # [(text, is_checked)]
        i = 0
        while i < len(text):
            if text[i] == self.CHECKBOX_CHECKED:
                segments.append((self.CHECKBOX_CHECKED, True))
                i += 1
            elif text[i] == self.CHECKBOX_UNCHECKED:
                segments.append((self.CHECKBOX_UNCHECKED, False))
                i += 1
            else:
                # 收集连续的非复选框字符
                start = i
                while i < len(text) and text[i] not in [self.CHECKBOX_CHECKED, self.CHECKBOX_UNCHECKED]:
                    i += 1
                segments.append((text[start:i], None))

        if len(segments) <= 1:
            return

        # 获取run在段落中的位置
        p_element = paragraph._p
        run_index = list(p_element).index(run_element)

        # 修改第一个segment
        first_text, first_is_checked = segments[0]
        run.text = first_text
        if first_is_checked is not None:
            run.font.color.theme_color = self.COLOR_SELECTED if first_is_checked else self.COLOR_UNSELECTED

        # 为其他segments创建新run
        NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
        for seg_text, is_checked in segments[1:]:
            new_r = etree.Element(f'{NS}r')
            new_rPr = etree.SubElement(new_r, f'{NS}rPr')

            # 复制原始的sym元素（符号字体设置）- 这是关键！
            if original_sym is not None:
                new_sym = etree.SubElement(new_rPr, f'{NS}sym')
                for attr in original_sym.attrib:
                    new_sym.set(attr, original_sym.get(attr))

            # 复制原始字号
            if original_sz:
                new_rPr.set('sz', original_sz)

            # 复制原始语言设置
            if original_lang:
                new_rPr.set('lang', original_lang)

            # 设置颜色
            if is_checked is not None:
                solidFill = etree.SubElement(new_rPr, f'{NS}solidFill')
                schemeClr = etree.SubElement(solidFill, f'{NS}schemeClr')
                if is_checked:
                    schemeClr.set('val', 'accent2')
                else:
                    schemeClr.set('val', 'accent1')

            # 设置文本
            t = etree.SubElement(new_r, f'{NS}t')
            t.text = seg_text

            # 处理空格
            if seg_text and (seg_text[0] == ' ' or seg_text[-1] == ' ' or '  ' in seg_text):
                t.set('{http://www.w3.org/XML/1998/namespace}space', 'preserve')

            # 插入到段落
            run_index += 1
            p_element.insert(run_index, new_r)

    def get_current_values(self) -> Dict[str, any]:
        result = {"evaluations": [], "overall": None, "homework": None}
        for _, shape in self._evaluation_shapes:
            result["evaluations"].append(self._extract_selected(shape.text, self.EVALUATION_OPTIONS))
        if self._overall_shape:
            result["overall"] = self._extract_selected(self._overall_shape.text, self.OVERALL_OPTIONS)
        if self._homework_shape:
            result["homework"] = self._extract_selected(self._homework_shape.text, self.HOMEWORK_OPTIONS)
        return result

    def _extract_selected(self, text: str, options: List[str]) -> Optional[str]:
        for opt in options:
            if f"{opt}{self.CHECKBOX_CHECKED}" in text:
                return opt
        return None

    def reset_all(self):
        for _, shape in self._evaluation_shapes:
            self._reset_shape(shape)
        if self._overall_shape:
            self._reset_shape(self._overall_shape)
        if self._homework_shape:
            self._reset_shape(self._homework_shape)

    def _reset_shape(self, shape):
        if not hasattr(shape, 'text_frame'):
            return
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = run.text.replace(self.CHECKBOX_CHECKED, self.CHECKBOX_UNCHECKED)
