# -*- coding: utf-8 -*-
"""
PPT生成器核心类
负责生成课程反馈PPT
"""

from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple
import io
import copy

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree

from .models import CourseUnitData, LayoutConfig
from .image_processor import ImageProcessor


class PPTGenerator:
    """PPT生成器类"""

    # 布局索引映射（从模板分析得出）
    LAYOUT_INDEX = {
        "cover": 0,           # 4_标题和内容 - 首页
        "course_info": 1,     # 标题和内容 - 课程信息页
        "work": 2,            # 3_标题和内容 - 精彩瞬间/作品页
        "vehicle": 3,         # 5_标题和内容 - 车辆展示页
        "program": 4,         # 2_标题和内容 - 程序展示页
        "double_image": 5,    # 1_标题和内容 - 精彩瞬间(双图)
        "model": 6,           # 6_标题和内容 - 模型展示页
    }

    # 布局名称映射
    LAYOUT_NAMES = {
        "cover": "4_标题和内容",
        "course_info": "标题和内容",
        "work": "3_标题和内容",
        "vehicle": "5_标题和内容",
        "program": "2_标题和内容",
        "double_image": "1_标题和内容",
        "model": "6_标题和内容",
    }

    # 图片位置配置（EMU单位）
    IMAGE_POSITIONS = {
        "model": {
            "left": 507918,
            "top": 1436402,
            "width": 5842163,
            "height": 6546959
        },
        "work": {
            "left": 405000,
            "top": 1114428,
            "width": 6048000,
            "height": 7560000
        },
        "vehicle": {
            "left": 405000,
            "top": 1114428,
            "width": 6048000,
            "height": 7560000
        },
        "program": {
            "left": 405000,
            "top": 1114428,
            "width": 6048000,
            "height": 7560000
        },
        "double_image": {
            "left1": 400000,
            "top1": 1114428,
            "width1": 3000000,
            "height1": 4000000,
            "left2": 3500000,
            "top2": 1114428,
            "width2": 3000000,
            "height2": 4000000
        }
    }

    def __init__(self, template_path: str):
        """
        初始化PPT生成器

        Args:
            template_path: 模板文件路径
        """
        self.template_path = Path(template_path)
        self.prs: Optional[Presentation] = None
        self._layout_cache: Dict[str, Any] = {}
        self.image_processor = ImageProcessor()
        self._load_template()

    def _load_template(self):
        """加载模板文件"""
        if not self.template_path.exists():
            raise FileNotFoundError(f"模板文件不存在: {self.template_path}")

        self.prs = Presentation(str(self.template_path))
        self._cache_layouts()

    def _cache_layouts(self):
        """缓存所有布局对象"""
        self._layout_cache = {}
        for key, index in self.LAYOUT_INDEX.items():
            if index < len(self.prs.slide_layouts):
                self._layout_cache[key] = self.prs.slide_layouts[index]

    def get_layout_by_name(self, name: str):
        """
        根据名称获取布局

        Args:
            name: 布局名称

        Returns:
            SlideLayout对象
        """
        for layout in self.prs.slide_layouts:
            if layout.name == name:
                return layout
        raise ValueError(f"找不到布局: {name}")

    def get_layout(self, layout_key: str):
        """
        根据键名获取布局

        Args:
            layout_key: 布局键名 (cover, course_info, work, vehicle, program, double_image, model)

        Returns:
            SlideLayout对象
        """
        if layout_key in self._layout_cache:
            return self._layout_cache[layout_key]

        # 尝试通过名称查找
        if layout_key in self.LAYOUT_NAMES:
            return self.get_layout_by_name(self.LAYOUT_NAMES[layout_key])

        raise ValueError(f"找不到布局: {layout_key}")

    def get_all_layouts(self) -> Dict[str, Any]:
        """
        获取所有布局

        Returns:
            布局字典 {key: layout}
        """
        return self._layout_cache.copy()

    def get_layout_info(self) -> List[Dict[str, Any]]:
        """
        获取布局信息列表

        Returns:
            布局信息列表
        """
        info = []
        for key, layout in self._layout_cache.items():
            info.append({
                "key": key,
                "name": layout.name,
                "index": self.LAYOUT_INDEX.get(key, -1)
            })
        return info

    def add_slide(self, layout_key: str):
        """
        添加新幻灯片

        Args:
            layout_key: 布局键名

        Returns:
            新添加的Slide对象
        """
        layout = self.get_layout(layout_key)
        slide = self.prs.slides.add_slide(layout)
        return slide

    def add_course_unit(self, unit_data: CourseUnitData, layout_config: LayoutConfig):
        """
        添加一个课程单元

        Args:
            unit_data: 课程单元数据
            layout_config: 母版配置
        """
        layout_sequence = self._get_layout_sequence(layout_config)

        for layout_key in layout_sequence:
            self.add_slide(layout_key)

    def _get_layout_sequence(self, config: LayoutConfig,
                              model_count: int = 1, work_count: int = 1,
                              program_count: int = 1, vehicle_count: int = 1) -> List[str]:
        """
        根据配置获取布局顺序

        Args:
            config: 布局配置
            model_count: 模型展示页数量
            work_count: 精彩瞬间页数量
            program_count: 程序展示页数量
            vehicle_count: 车辆展示页数量

        Returns:
            布局键名列表
        """
        sequence = []

        if config.include_course_info:
            sequence.append("course_info")

        if config.include_model_display:
            for _ in range(model_count):
                sequence.append("model")

        if config.include_work_display:
            for _ in range(work_count):
                sequence.append("work")

        if config.include_program_display:
            for _ in range(program_count):
                sequence.append("program")

        if config.include_vehicle_display:
            for _ in range(vehicle_count):
                sequence.append("vehicle")

        return sequence

    def clear_slides(self):
        """清除所有幻灯片（保留模板）"""
        # 获取幻灯片数量
        slide_count = len(self.prs.slides)
        # 从后往前删除
        rId_list = []
        for slide in self.prs.slides:
            rId = self.prs.part.relate_to(slide.part, 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide')
            rId_list.append(rId)

        # 注意：python-pptx不支持直接删除幻灯片，需要操作XML
        # 这里采用重新加载模板的方式
        self._load_template()

    def create_new_presentation(self):
        """基于模板创建新的演示文稿（清除原有内容）"""
        # 保存模板路径
        template_path = self.template_path
        # 重新加载
        self.prs = Presentation(str(template_path))
        self._cache_layouts()

    def save(self, output_path: str) -> bool:
        """
        保存PPT文件

        Args:
            output_path: 输出文件路径

        Returns:
            是否保存成功
        """
        try:
            # 确保输出目录存在
            output = Path(output_path)
            output.parent.mkdir(parents=True, exist_ok=True)

            self.prs.save(str(output))
            return True
        except Exception as e:
            print(f"保存失败: {e}")
            return False

    def save_as(self, output_path: str) -> bool:
        """
        另存为PPT文件

        Args:
            output_path: 输出文件路径

        Returns:
            是否保存成功
        """
        return self.save(output_path)

    def get_slide_count(self) -> int:
        """获取当前幻灯片数量"""
        return len(self.prs.slides)

    def get_template_path(self) -> str:
        """获取模板路径"""
        return str(self.template_path)

    def verify_template(self) -> Dict[str, Any]:
        """
        验证模板是否有效

        Returns:
            验证结果字典
        """
        result = {
            "valid": True,
            "template_path": str(self.template_path),
            "layout_count": len(self.prs.slide_layouts),
            "layouts": {},
            "errors": []
        }

        # 检查所有需要的布局是否存在
        for key, expected_name in self.LAYOUT_NAMES.items():
            found = False
            for layout in self.prs.slide_layouts:
                if layout.name == expected_name:
                    found = True
                    result["layouts"][key] = {
                        "name": layout.name,
                        "found": True
                    }
                    break

            if not found:
                result["layouts"][key] = {
                    "name": expected_name,
                    "found": False
                }
                result["errors"].append(f"缺少布局: {expected_name}")
                result["valid"] = False

        return result

    # ==================== 图片嵌入功能 ====================

    def find_picture_shape(self, slide) -> Optional[Any]:
        """
        查找幻灯片中的图片形状

        Args:
            slide: Slide对象

        Returns:
            找到的Picture形状，如果没有返回None
        """
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return shape
        return None

    def remove_picture_shape(self, slide) -> bool:
        """
        移除幻灯片中的图片形状

        Args:
            slide: Slide对象

        Returns:
            是否成功移除
        """
        shape_to_remove = self.find_picture_shape(slide)
        if shape_to_remove:
            sp = shape_to_remove._element
            sp.getparent().remove(sp)
            return True
        return False

    def insert_image(
        self,
        slide,
        image_path: str,
        left: float = None,
        top: float = None,
        width: float = None,
        height: float = None,
        layout_type: str = None
    ) -> bool:
        """
        在幻灯片中插入图片

        Args:
            slide: Slide对象
            image_path: 图片路径
            left: 左边距（EMU）
            top: 上边距（EMU）
            width: 宽度（EMU）
            height: 高度（EMU）
            layout_type: 布局类型，如果指定则自动裁剪保持比例

        Returns:
            是否成功
        """
        try:
            # 检查图片文件是否存在
            if not Path(image_path).exists():
                print(f"图片文件不存在: {image_path}")
                return False

            # 先移除现有图片
            self.remove_picture_shape(slide)

            # 使用默认位置（如果未指定）
            if left is None:
                left = 400000
            if top is None:
                top = 1100000
            if width is None:
                width = 6000000
            if height is None:
                height = 7500000

            # 如果指定了布局类型，先处理图片（保持比例裁剪）
            if layout_type:
                processed_img = self.image_processor.process_image_for_ppt(image_path, layout_type)
                if processed_img:
                    # 将处理后的图片保存到字节流
                    img_bytes = io.BytesIO()
                    processed_img.save(img_bytes, format='PNG')
                    img_bytes.seek(0)
                    # 使用字节流插入图片
                    slide.shapes.add_picture(img_bytes, left, top, width, height)
                    return True
                else:
                    print(f"图片处理失败，使用原始图片: {image_path}")

            # 直接添加图片（不处理）
            slide.shapes.add_picture(image_path, left, top, width, height)
            return True

        except Exception as e:
            print(f"插入图片失败: {e}")
            return False

    def insert_image_from_bytes(
        self,
        slide,
        image_bytes: bytes,
        left: float,
        top: float,
        width: float = None,
        height: float = None
    ) -> bool:
        """
        从字节数据插入图片

        Args:
            slide: Slide对象
            image_bytes: 图片字节数据
            left: 左边距（EMU）
            top: 上边距（EMU）
            width: 宽度（EMU）
            height: 高度（EMU）

        Returns:
            是否成功
        """
        try:
            # 先移除现有图片
            self.remove_picture_shape(slide)

            # 创建字节流
            image_stream = io.BytesIO(image_bytes)

            # 添加图片
            if width and height:
                slide.shapes.add_picture(image_stream, left, top, width, height)
            elif width:
                slide.shapes.add_picture(image_stream, left, top, width=width)
            else:
                slide.shapes.add_picture(image_stream, left, top)

            return True

        except Exception as e:
            print(f"从字节插入图片失败: {e}")
            return False

    def insert_image_for_layout(
        self,
        slide,
        image_path: str,
        layout_type: str
    ) -> bool:
        """
        根据布局类型插入图片到指定位置

        Args:
            slide: Slide对象
            image_path: 图片路径
            layout_type: 布局类型 (model, work, vehicle, program, double_image)

        Returns:
            是否成功
        """
        if layout_type not in self.IMAGE_POSITIONS:
            print(f"未知的布局类型: {layout_type}")
            return False

        pos = self.IMAGE_POSITIONS[layout_type]

        # 双图布局特殊处理
        if layout_type == "double_image":
            return self.insert_image(
                slide, image_path,
                pos["left1"], pos["top1"], pos["width1"], pos["height1"],
                layout_type="double_image"
            )

        return self.insert_image(
            slide, image_path,
            pos["left"], pos["top"], pos["width"], pos["height"],
            layout_type=layout_type
        )

    def insert_double_images(
        self,
        slide,
        image_path1: str,
        image_path2: str
    ) -> Tuple[bool, bool]:
        """
        在双图布局中插入两张图片

        Args:
            slide: Slide对象
            image_path1: 第一张图片路径
            image_path2: 第二张图片路径

        Returns:
            (第一张是否成功, 第二张是否成功)
        """
        pos = self.IMAGE_POSITIONS["double_image"]

        # 移除现有图片
        self.remove_picture_shape(slide)

        result1 = self.insert_image(
            slide, image_path1,
            pos["left1"], pos["top1"], pos["width1"], pos["height1"],
            layout_type="double_image"
        )

        result2 = self.insert_image(
            slide, image_path2,
            pos["left2"], pos["top2"], pos["width2"], pos["height2"],
            layout_type="double_image"
        )

        return (result1, result2)

    def get_image_position(self, layout_type: str) -> Optional[Dict[str, float]]:
        """
        获取指定布局的图片位置信息

        Args:
            layout_type: 布局类型

        Returns:
            位置信息字典
        """
        return self.IMAGE_POSITIONS.get(layout_type)

    def add_slide_with_image(
        self,
        layout_key: str,
        image_path: str
    ):
        """
        添加带图片的幻灯片

        Args:
            layout_key: 布局键名
            image_path: 图片路径

        Returns:
            新添加的Slide对象，失败返回None
        """
        slide = self.add_slide(layout_key)
        if slide:
            success = self.insert_image_for_layout(slide, image_path, layout_key)
            if success:
                return slide

        return None

    # ==================== 幻灯片管理功能 (F011) ====================

    def duplicate_slide(self, slide_index: int) -> Optional[Any]:
        """
        复制指定索引的幻灯片（包括所有内容）

        Args:
            slide_index: 要复制的幻灯片索引

        Returns:
            新幻灯片对象，失败返回None
        """
        try:
            if slide_index < 0 or slide_index >= len(self.prs.slides):
                return None

            source_slide = self.prs.slides[slide_index]

            # 获取源幻灯片的布局
            source_layout = source_slide.slide_layout

            # 添加新幻灯片
            new_slide = self.prs.slides.add_slide(source_layout)

            # 复制幻灯片内容
            self._copy_slide_content(source_slide, new_slide)

            return new_slide

        except Exception as e:
            print(f"复制幻灯片失败: {e}")
            return None

    def _copy_slide_content(self, source_slide, target_slide):
        """
        复制幻灯片内容（形状、文本等，但不复制图片）

        注意：不复制图片，因为图片的引用ID会导致冲突。
        图片会在后续的填充流程中重新插入。

        Args:
            source_slide: 源幻灯片
            target_slide: 目标幻灯片
        """
        # 获取源幻灯片的形状列表
        source_shapes = source_slide.shapes

        # 复制每个形状（跳过图片）
        for shape in source_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 跳过图片复制，避免引用ID冲突
                pass
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # 复制组合形状
                self._copy_group_shape(shape, target_slide)
            else:
                # 复制其他形状
                self._copy_shape(shape, target_slide)

    def _copy_shape(self, source_shape, target_slide):
        """复制普通形状（清理图片引用以避免冲突）"""
        try:
            # 使用XML深复制
            el = source_shape._element
            new_el = copy.deepcopy(el)

            # 清理图片引用（避免rId冲突）
            for blip in new_el.iter(qn('a:blip')):
                if 'r:embed' in blip.attrib:
                    del blip.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed']

            # 添加到目标幻灯片
            target_slide.shapes._spTree.insert_element_before(
                new_el,
                'p:extLst'
            )
        except Exception as e:
            pass

    def _copy_picture_shape(self, source_shape, target_slide):
        """复制图片形状"""
        try:
            # 使用XML深复制
            el = source_shape._element
            new_el = copy.deepcopy(el)

            # 添加到目标幻灯片
            target_slide.shapes._spTree.insert_element_before(
                new_el,
                'p:extLst'
            )
        except Exception as e:
            pass

    def _copy_group_shape(self, source_shape, target_slide):
        """复制组合形状（清理图片引用以避免冲突）"""
        try:
            # 使用XML深复制
            el = source_shape._element
            new_el = copy.deepcopy(el)

            # 清理图片引用（避免rId冲突）
            # 在组合形状中查找所有 a:blip 元素并移除 r:embed 属性
            for blip in new_el.iter(qn('a:blip')):
                if 'r:embed' in blip.attrib:
                    del blip.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed']

            # 添加到目标幻灯片
            target_slide.shapes._spTree.insert_element_before(
                new_el,
                'p:extLst'
            )
        except Exception as e:
            pass

    def delete_slide(self, slide_index: int) -> bool:
        """
        删除指定索引的幻灯片

        Args:
            slide_index: 要删除的幻灯片索引

        Returns:
            是否成功删除
        """
        try:
            if slide_index < 0 or slide_index >= len(self.prs.slides):
                return False

            slide = self.prs.slides[slide_index]

            # 获取幻灯片的关系ID
            rId = self.prs.part.relate_to(
                slide.part,
                'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
            )

            # 从幻灯片ID列表中移除
            sldIdLst = self.prs.part._element.get_or_add_sldIdLst()
            for sldId in sldIdLst:
                if sldId.rId == rId:
                    sldIdLst.remove(sldId)
                    break

            return True

        except Exception as e:
            print(f"删除幻灯片失败: {e}")
            return False

    def delete_slides_after(self, keep_count: int) -> int:
        """
        删除指定数量之后的幻灯片

        Args:
            keep_count: 要保留的幻灯片数量

        Returns:
            删除的幻灯片数量
        """
        deleted = 0
        while len(self.prs.slides) > keep_count:
            if self.delete_slide(len(self.prs.slides) - 1):
                deleted += 1
            else:
                break
        return deleted

    def get_slide_by_index(self, index: int) -> Optional[Any]:
        """
        获取指定索引的幻灯片

        Args:
            index: 幻灯片索引

        Returns:
            Slide对象，不存在返回None
        """
        if 0 <= index < len(self.prs.slides):
            return self.prs.slides[index]
        return None

    def get_slide_count(self) -> int:
        """获取当前幻灯片数量"""
        return len(self.prs.slides)

    # ==================== 课程单元生成功能 ====================

    # 模板中幻灯片布局的索引映射
    TEMPLATE_SLIDE_INDEX = {
        "cover": 0,          # 封面页
        "course_info": 1,    # 课程信息页
        "model": 2,          # 模型展示页
        "work": 3,           # 作品展示页
    }

    def generate_from_template(self, layout_config: LayoutConfig,
                                model_image_count: int = 0,
                                work_image_count: int = 0,
                                program_image_count: int = 0,
                                vehicle_image_count: int = 0,
                                include_cover: bool = True) -> bool:
        """
        根据配置从模板生成PPT

        根据include_cover参数决定是否保留封面页。
        支持根据实际上传的图片数量动态生成对应数量的幻灯片。

        重要：采用"先添加后删除"的顺序，避免 python-pptx 的 bug。
        python-pptx 在删除幻灯片后添加新幻灯片会产生重复的 ZIP 条目，
        导致生成的 PPT 文件损坏。

        页面顺序：
        - 封面 (索引0，仅当include_cover=True)
        - 课程信息 (索引0或1，取决于是否有封面)
        - 模型展示页 (连续N页，每页一张图)
        - 程序展示页 (连续P页，每页一张图)
        - 车辆展示页 (连续V页，每页一张图)
        - 精彩瞬间 (连续M页，每页一张图)

        Args:
            layout_config: 布局配置
            model_image_count: 模型展示图片数量（决定模型展示页数量）
            work_image_count: 精彩瞬间图片数量
            program_image_count: 程序展示图片数量
            vehicle_image_count: 车辆展示图片数量
            include_cover: 是否包含封面页（第1课为True，其他为False）

        Returns:
            是否成功
        """
        try:
            import tempfile
            import os

            # 计算各类型页面数量
            model_count = model_image_count if layout_config.include_model_display else 0
            work_count = work_image_count if layout_config.include_work_display else 0
            program_count = program_image_count if layout_config.include_program_display else 0
            vehicle_count = vehicle_image_count if layout_config.include_vehicle_display else 0

            # 总页数 = 封面(可选) + 课程信息(1) + 各类型页面
            base_pages = 2 if include_cover else 1  # 封面+课程信息 或 仅课程信息
            total_slides = base_pages + model_count + work_count + program_count + vehicle_count
            new_slides_count = model_count + work_count + program_count + vehicle_count

            # 重新加载模板
            self._load_template()

            # 步骤1：获取布局
            model_layout = self.get_layout("model")
            work_layout = self.get_layout("work")
            program_layout = self.get_layout("program")
            vehicle_layout = self.get_layout("vehicle")

            # 步骤2：【先添加】新幻灯片（关键：必须在删除之前添加！）
            # 顺序：模型展示 -> 程序展示 -> 车辆展示 -> 精彩瞬间
            for i in range(model_count):
                self.prs.slides.add_slide(model_layout)

            for i in range(program_count):
                self.prs.slides.add_slide(program_layout)

            for i in range(vehicle_count):
                self.prs.slides.add_slide(vehicle_layout)

            for i in range(work_count):
                self.prs.slides.add_slide(work_layout)

            # 步骤3：【后删除】不需要的幻灯片
            # 当前结构：[原始79页] + [新添加N页]
            #
            # 需要保留的幻灯片：
            # - include_cover=True: 索引 0(封面), 1(课程信息), 79..(新添加)
            # - include_cover=False: 索引 1(课程信息), 79..(新添加)
            #
            # 需要删除的幻灯片：原始模板中不需要的页面

            original_count = len(self.prs.slides) - new_slides_count

            # 确定要保留的原始幻灯片索引
            if include_cover:
                # 保留索引 0(封面) 和 1(课程信息)
                keep_original_indices = {0, 1}
            else:
                # 只保留索引 1(课程信息)，不保留封面
                keep_original_indices = {1}

            # 从后往前删除原始幻灯片中不需要的
            # 索引范围: 0 到 original_count-1
            for idx in range(original_count - 1, -1, -1):
                if idx not in keep_original_indices:
                    rId = self.prs.slides._sldIdLst[idx].rId
                    self.prs.part.drop_rel(rId)
                    del self.prs.slides._sldIdLst[idx]

            # 步骤4：保存到临时文件
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
                tmp_path = tmp.name

            try:
                self.prs.save(tmp_path)

                # 注意：不再需要 _clean_pptx_file
                # 由于采用"先添加后删除"的顺序，python-pptx 不会产生重复的 ZIP 条目
                # 孤立的幻灯片文件会被保留在 PPT 中，但不影响使用
                # 直接重新加载保存的文件即可

                # 重新加载文件
                self.prs = Presentation(tmp_path)
                self._cache_layouts()

                # 删除临时文件
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)

            except Exception as e:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
                raise e

            return True

        except Exception as e:
            print(f"生成PPT失败: {e}")
            import traceback
            traceback.print_exc()
            return False

    def _clean_pptx_file(self, file_path: str, keep_slide_count: int) -> str:
        """
        清理PPTX文件中多余的幻灯片文件和冗余关系

        python-pptx的delete_slide只删除sldIdLst中的引用，但不会：
        1. 删除实际的幻灯片文件
        2. 删除presentation.xml.rels中的关系条目

        此外，python-pptx在某些操作后可能产生重复映射（同一文件被多个rId引用）。

        此方法会：
        1. 清理重复的rId映射（同一文件只保留一个rId）
        2. 清理presentation.xml.rels中不再需要的幻灯片关系
        3. 更新presentation.xml中的sldIdLst
        4. 删除孤立的幻灯片文件

        Args:
            file_path: 原始PPTX文件路径
            keep_slide_count: 要保留的幻灯片数量

        Returns:
            清理后的文件路径
        """
        import zipfile
        import tempfile
        import os
        import xml.etree.ElementTree as ET

        # 创建临时文件
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            output_path = tmp.name

        with zipfile.ZipFile(file_path, 'r') as zin:
            # 定义命名空间
            REL_NS = '{http://schemas.openxmlformats.org/package/2006/relationships}'
            SLIDE_REL_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
            P_NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            R_NS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

            # 1. 读取 presentation.xml.rels 获取所有幻灯片关系
            rels_content = zin.read('ppt/_rels/presentation.xml.rels')
            rels_root = ET.fromstring(rels_content)

            # 构建 rId -> 文件名 的映射
            all_rid_to_file = {}
            for rel in rels_root.findall(f'.//{REL_NS}Relationship'):
                rid = rel.get('Id')
                target = rel.get('Target')
                rel_type = rel.get('Type')
                if rel_type == SLIDE_REL_TYPE:
                    all_rid_to_file[rid] = 'ppt/' + target

            # 2. 读取 presentation.xml 获取 sldIdLst
            pres_content = zin.read('ppt/presentation.xml')
            pres_root = ET.fromstring(pres_content)

            # 找到 sldIdLst 元素
            sldIdLst_elem = pres_root.find(f'.//{{{P_NS}}}sldIdLst')
            if sldIdLst_elem is None:
                # 如果没有 sldIdLst，直接返回原文件
                return file_path

            # 获取前 keep_slide_count 个 rId（保持顺序）
            sldId_elems = list(sldIdLst_elem)
            original_rids = []
            for i, sldId_elem in enumerate(sldId_elems):
                if i < keep_slide_count:
                    rid = sldId_elem.get(f'{{{R_NS}}}id')
                    original_rids.append(rid)

            # 3. 处理重复映射：同一文件只保留第一个 rId
            file_to_rid = {}  # 文件 -> 保留的 rId
            rid_to_keep = set()  # 最终要保留的 rId 集合
            rids_to_remove_from_sldIdLst = set()  # 需要从 sldIdLst 中移除的 rId

            for rid in original_rids:
                if rid in all_rid_to_file:
                    file_path_name = all_rid_to_file[rid]
                    if file_path_name not in file_to_rid:
                        # 第一次见到这个文件，保留这个 rId
                        file_to_rid[file_path_name] = rid
                        rid_to_keep.add(rid)
                    else:
                        # 这个文件已经被另一个 rId 引用，标记此 rId 为重复
                        rids_to_remove_from_sldIdLst.add(rid)

            # 4. 清理 presentation.xml.rels：删除不在 rid_to_keep 中的幻灯片关系
            for rel in list(rels_root.findall(f'.//{REL_NS}Relationship')):
                rid = rel.get('Id')
                rel_type = rel.get('Type')
                if rel_type == SLIDE_REL_TYPE and rid not in rid_to_keep:
                    rels_root.remove(rel)

            cleaned_rels_content = ET.tostring(rels_root, encoding='UTF-8', xml_declaration=True)

            # 5. 清理 presentation.xml 中的 sldIdLst
            for sldId_elem in list(sldIdLst_elem):
                rid = sldId_elem.get(f'{{{R_NS}}}id')
                if rid in rids_to_remove_from_sldIdLst or rid not in rid_to_keep:
                    sldIdLst_elem.remove(sldId_elem)

            cleaned_pres_content = ET.tostring(pres_root, encoding='UTF-8', xml_declaration=True)

            # 6. 构建要保留的幻灯片文件集合
            keep_files = set()
            for rid in rid_to_keep:
                if rid in all_rid_to_file:
                    slide_file = all_rid_to_file[rid]
                    keep_files.add(slide_file)
                    # 也保留对应的 .rels 文件
                    slide_num = slide_file.split('slide')[1].split('.')[0]
                    rels_file = f'ppt/slides/_rels/slide{slide_num}.xml.rels'
                    keep_files.add(rels_file)

            # 7. 复制文件，清理不需要的
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zout:
                for item in zin.namelist():
                    if item == 'ppt/_rels/presentation.xml.rels':
                        # 写入清理后的关系文件
                        zout.writestr(item, cleaned_rels_content)
                    elif item == 'ppt/presentation.xml':
                        # 写入清理后的 presentation.xml
                        zout.writestr(item, cleaned_pres_content)
                    elif item.startswith('ppt/slides/slide') and item.endswith('.xml'):
                        # 只保留需要的幻灯片文件
                        if item in keep_files:
                            zout.writestr(item, zin.read(item))
                    elif item.startswith('ppt/slides/_rels/slide') and item.endswith('.rels'):
                        # 只保留需要的幻灯片关系文件
                        if item in keep_files:
                            zout.writestr(item, zin.read(item))
                    else:
                        # 保留其他所有文件
                        zout.writestr(item, zin.read(item))

        return output_path

    def _delete_slides_after_index(self, keep_count: int):
        """
        删除指定索引之后的所有幻灯片

        Args:
            keep_count: 要保留的幻灯片数量
        """
        while len(self.prs.slides) > keep_count:
            slide_index = len(self.prs.slides) - 1
            self.delete_slide(slide_index)

    def clear_content_slides(self):
        """清除内容幻灯片，只保留封面"""
        self.delete_slides_after(1)

    # ==================== PPT保存功能 (F012) ====================

    @staticmethod
    def generate_filename(
        student_name: str,
        lesson_number: int = None,
        prefix: str = "课程反馈",
        date_format: str = "%Y%m%d_%H%M%S"
    ) -> str:
        """
        生成PPT文件名

        Args:
            student_name: 学生姓名
            lesson_number: 课时编号（可选）
            prefix: 文件名前缀
            date_format: 日期格式

        Returns:
            生成的文件名（不含路径）
        """
        from datetime import datetime

        # 构建文件名
        parts = [student_name]

        if lesson_number:
            parts.append(f"第{lesson_number}课")

        parts.append(prefix)
        parts.append(datetime.now().strftime(date_format))

        filename = "_".join(parts) + ".pptx"
        return filename

    @staticmethod
    def generate_pdf_filename(ppt_filename: str) -> str:
        """
        从PPT文件名生成对应的PDF文件名（去除时间戳）

        Args:
            ppt_filename: PPT文件名，如 "亮亮_第4课_课程反馈_20260307_082001.pptx"

        Returns:
            PDF文件名，如 "亮亮_第4课_课程反馈.pdf"
        """
        import re
        # 移除时间戳后缀 _YYYYMMDD_HHMMSS.pptx
        pdf_name = re.sub(r'_\d{8}_\d{6}\.pptx$', '.pdf', ppt_filename)
        return pdf_name

    def save_with_data(
        self,
        output_dir: str,
        student_name: str,
        lesson_number: int = None,
        prefix: str = "课程反馈",
        overwrite: bool = True
    ) -> Tuple[bool, str]:
        """
        根据课程数据保存PPT

        Args:
            output_dir: 输出目录
            student_name: 学生姓名
            lesson_number: 课时编号
            prefix: 文件名前缀
            overwrite: 是否覆盖已存在的文件

        Returns:
            (是否成功, 实际保存路径)
        """
        # 生成文件名
        filename = self.generate_filename(student_name, lesson_number, prefix)
        output_path = Path(output_dir) / filename

        # 检查文件是否存在
        if output_path.exists() and not overwrite:
            # 生成新文件名（添加序号）
            counter = 1
            while output_path.exists():
                stem = output_path.stem
                new_filename = f"{stem}_{counter}.pptx"
                output_path = Path(output_dir) / new_filename
                counter += 1

        # 保存
        success = self.save(str(output_path))
        return (success, str(output_path))

    def save_course_unit(
        self,
        output_dir: str,
        unit_data: 'CourseUnitData',
        overwrite: bool = True
    ) -> Tuple[bool, str]:
        """
        保存课程单元PPT

        Args:
            output_dir: 输出目录
            unit_data: 课程单元数据
            overwrite: 是否覆盖已存在的文件

        Returns:
            (是否成功, 实际保存路径)
        """
        return self.save_with_data(
            output_dir,
            unit_data.student_name,
            unit_data.lesson_number,
            "课程反馈",
            overwrite
        )
