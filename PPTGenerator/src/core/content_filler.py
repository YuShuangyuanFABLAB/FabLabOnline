# -*- coding: utf-8 -*-
"""
PPT内容填充器
将课程数据填充到PPT模板中
"""

from typing import List
from pathlib import Path
import copy

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from src.core.models import CourseUnitData
from src.core.checkbox_handler import CheckboxHandler

# 默认文字颜色
DEFAULT_TEXT_COLOR = RGBColor(0x40, 0x40, 0x40)  # #404040
HIGHLIGHT_COLOR = RGBColor(0x00, 0x70, 0xC0)  # 重点-蓝色
DIFFICULTY_COLOR = RGBColor(0xED, 0x7D, 0x31)  # 难点-橙色

# 模板默认字体（与模板保持一致）
TEMPLATE_DEFAULT_FONT = "-apple-system"  # 西文字体
TEMPLATE_DEFAULT_FONT_EAST = "等线"  # 中文字体（东亚字体）


class ContentFiller:
    """PPT内容填充器"""

    def __init__(self, presentation: Presentation):
        self.prs = presentation

    def fill_course_info(self, data: CourseUnitData, series_name: str = None, series_level: int = None,
                         slide_index: int = 1) -> bool:
        """
        填充课程信息页

        Args:
            data: 课程数据
            series_name: 课程系列名称（如"机械臂设计"）
            series_level: 课程阶数（如1）
            slide_index: 课程信息页的索引（有封面时为1，无封面时为0）

        Returns:
            是否成功
        """
        if len(self.prs.slides) < slide_index + 1:
            return False

        slide = self.prs.slides[slide_index]

        # 构建系列名称
        if series_name and series_level:
            pass  # series_full computed for documentation purposes
        else:
            series_name = "机械臂设计"

        # ==================== 修改幻灯片上的内容 ====================
        # 1. 替换基本信息（直接搜索并替换模板中的值）
        # 课时编号（使用白色字体）
        self._replace_lesson_number(slide, '（第 1 课）', f'（第 {data.lesson_number} 课）')

        # 课程内容（替换幻灯片上的"机械臂设计（1阶）"为课程内容）
        # 如果用户未填写，则使用"系列名(x阶)"格式
        if data.course_content and data.course_content.strip():
            course_content_text = data.course_content
        elif series_name and series_level:
            course_content_text = f"{series_name}（{series_level}阶）"
        else:
            course_content_text = data.course_content  # 使用原值（可能为空）
        self._replace_text_in_slide(slide, '机械臂设计（1阶）', course_content_text)

        # 学生姓名（直接替换模板中的学生名）
        self._replace_text_in_slide(slide, '赵如一', data.student_name)

        # 授课教师
        self._replace_text_in_slide(slide, '于双源', data.teacher_name)

        # 课时数
        self._replace_text_in_slide(slide, '2课时', f'{data.class_hours}课时')

        # 上课日期
        self._replace_text_in_slide(slide, '2025.7.23 18:30-20:30', data.class_date)

        # 2. 处理评价复选框
        handler = CheckboxHandler(slide)

        # PPT模板中评价项的顺序（先左列从上到下，再右列从上到下）：
        # 左列：
        # 0. 专注度与效率 (focus_efficiency)
        # 1. 课堂听课习惯 (listening_habit)
        # 2. 课堂任务完成 (task_completion)
        # 3. 课堂内容理解 (content_understanding)
        # 4. 知识熟练程度 (knowledge_proficiency)
        # 5. 学习方法习惯 (learning_method)
        # 右列：
        # 6. 动手实践能力 (hands_on_ability)
        # 7. 逻辑思维表现 (logic_thinking)
        # 8. 想象力与创新 (imagination_creativity)
        # 9. 独立分析思考 (independent_analysis)
        # 10. 解决问题表现 (problem_solving)
        # 11. 挫折困难应对 (frustration_handling)
        evaluations = [
            data.focus_efficiency,        # 0 - 左列第1个
            data.listening_habit,         # 1 - 左列第2个
            data.task_completion,         # 2 - 左列第3个
            data.content_understanding,   # 3 - 左列第4个
            data.knowledge_proficiency,   # 4 - 左列第5个
            data.learning_method,         # 5 - 左列第6个
            data.hands_on_ability,        # 6 - 右列第1个
            data.logic_thinking,          # 7 - 右列第2个
            data.imagination_creativity,  # 8 - 右列第3个
            data.independent_analysis,    # 9 - 右列第4个
            data.problem_solving,         # 10 - 右列第5个
            data.frustration_handling,    # 11 - 右列第6个
        ]

        for i, level in enumerate(evaluations):
            if level and i < handler.get_evaluation_count():
                handler.set_evaluation(i, level.value)

        if data.overall_evaluation:
            handler.set_overall_evaluation(data.overall_evaluation.value)

        if data.last_homework_status:
            handler.set_homework_evaluation(data.last_homework_status.value)

        # 3. 替换补充说明等文本（找到包含关键词的形状，整体替换）
        if data.additional_comments:
            self._replace_shape_containing(slide, '今天开始正式学习', data.additional_comments)
            self._replace_shape_containing(slide, '补充说明', data.additional_comments, is_label=True)

        if data.homework:
            self._replace_shape_containing(slide, '本次课无课堂作业', data.homework)
            self._replace_shape_containing(slide, '课程作业', data.homework, is_label=True)

        if data.other_notes:
            self._replace_shape_containing(slide, '三维建模对空间想象力', data.other_notes)
            self._replace_shape_containing(slide, '注意事项', data.other_notes, is_label=True)

        # 4. 处理课堂知识内容（带颜色标记）
        if data.knowledge_content:
            self._fill_knowledge_content(slide, data.knowledge_content, data.highlights, data.difficulties)

        return True

    def fill_cover_series(self, series_name: str, series_level: int) -> bool:
        """
        填充封面页（首页）的课程系列名称

        模板中有两处需要替换（都在布局中）：
        1. 首页上方：机械臂设计（橙色 #FF6224）- 布局0中的文本框 614
        2. 首页下方：FABLAB法贝实验室 — — 机械臂设计（1阶）课程（浅蓝色 #D3F1F0）- 布局0中的文本框 750

        Args:
            series_name: 系列名称（如"机械臂设计"）
            series_level: 阶数（如1）

        Returns:
            是否成功
        """
        if len(self.prs.slide_masters) < 1:
            return False

        master = self.prs.slide_masters[0]
        layouts = master.slide_layouts

        # 封面页使用布局0 (4_标题和内容)
        if len(layouts) < 1:
            return False

        layout = layouts[0]  # 布局0
        series_full = f"{series_name}（{series_level}阶）"

        # 在布局中查找并替换
        # 1. 首页下方：FABLAB法贝实验室 — — 机械臂设计（1阶）课程
        self._replace_text_in_layout(layout, '机械臂设计（1阶）课程', series_full + '课程')

        # 2. 首页上方：机械臂设计
        self._replace_text_in_layout(layout, '机械臂设计', series_name)

        return True

    def _replace_text_in_layout(self, layout, old_text: str, new_text: str) -> bool:
        """
        在布局中查找并替换文本，保留所有格式

        Args:
            layout: 幻灯片布局
            old_text: 要替换的文本
            new_text: 新文本

        Returns:
            是否成功
        """
        shape = self._find_shape_in_layout(layout, old_text)
        if shape:
            return self._replace_layout_text_keep_format(shape, old_text, new_text)
        return False

    def _replace_layout_text_keep_format(self, shape, old_text: str, new_text: str) -> bool:
        """
        替换布局中的文本并完整保留格式（颜色、字体、对齐等）
        通过直接操作XML来确保格式不丢失

        Args:
            shape: 形状对象
            old_text: 要替换的文本
            new_text: 新文本

        Returns:
            是否成功
        """
        if not hasattr(shape, 'text_frame'):
            return False

        tf = shape.text_frame
        if not tf.paragraphs:
            return False

        # 检查形状的完整文本是否包含目标文本
        shape_text = shape.text
        if old_text not in shape_text:
            return False

        # 情况1：单个run包含完整的目标文本（简单情况）
        found_in_single_run = False
        for para in tf.paragraphs:
            for run in para.runs:
                if old_text in run.text:
                    t_elem = run._r.find(qn('a:t'))
                    if t_elem is not None:
                        if run.text == old_text:
                            t_elem.text = new_text
                        else:
                            t_elem.text = run.text.replace(old_text, new_text)
                        found_in_single_run = True

        if found_in_single_run:
            return True

        # 情况2：目标文本被分散在多个run中
        # 需要找到包含目标文本的run范围，然后合并替换
        for para in tf.paragraphs:
            runs = list(para.runs)
            if not runs:
                continue

            # 收集所有run的文本
            run_texts = [run.text for run in runs]
            full_text = "".join(run_texts)

            if old_text not in full_text:
                continue

            # 找到目标文本在完整文本中的位置
            start_pos = full_text.find(old_text)
            end_pos = start_pos + len(old_text)

            # 替换文本
            new_full_text = full_text[:start_pos] + new_text + full_text[end_pos:]

            # 获取第一个run的格式（用于保持格式）
            first_run = runs[0]
            rPr_elem = first_run._r.find(qn('a:rPr'))

            # 深拷贝第一个run的rPr
            rPr_copy = None
            if rPr_elem is not None:
                rPr_copy = copy.deepcopy(rPr_elem)

            # 删除所有run的XML元素
            p_elem = para._p
            for run in runs:
                run._r.getparent().remove(run._r)

            # 在endParaRPr之前插入新的run
            new_r = etree.Element(qn('a:r'))
            if rPr_copy is not None:
                new_r.append(rPr_copy)
            t_elem = etree.SubElement(new_r, qn('a:t'))
            t_elem.text = new_full_text

            # 找到endParaRPr元素，在其之前插入
            endParaRPr = p_elem.find(qn('a:endParaRPr'))
            if endParaRPr is not None:
                p_elem.insert(list(p_elem).index(endParaRPr), new_r)
            else:
                p_elem.append(new_r)

            return True

        return False

    def _get_run_color(self, run) -> str:
        """从run的XML中获取颜色值"""
        rPr = run._r.find(qn('a:rPr'))
        if rPr is not None:
            solidFill = rPr.find(qn('a:solidFill'))
            if solidFill is not None:
                srgbClr = solidFill.find(qn('a:srgbClr'))
                if srgbClr is not None:
                    return srgbClr.get('val')
        return None

    def _set_run_color(self, run, color_hex: str):
        """设置run的颜色（通过XML）"""
        from pptx.dml.color import RGBColor
        rPr = run._r.get_or_add_rPr()

        # 移除现有的填充
        for child in list(rPr):
            if child.tag == qn('a:solidFill'):
                rPr.remove(child)

        # 添加新的实心填充
        solidFill = etree.SubElement(rPr, qn('a:solidFill'))
        srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
        srgbClr.set('val', color_hex)

        # 同时设置python-pptx的颜色属性
        run.font.color.rgb = RGBColor.from_string(color_hex)

    def _find_shape_in_layout(self, layout, text: str):
        """
        在布局中查找包含指定文本的形状

        Args:
            layout: 幻灯片布局
            text: 要查找的文本

        Returns:
            找到的形状，或None
        """
        for shape in layout.shapes:
            if hasattr(shape, 'text') and text in shape.text:
                return shape
            # 检查组合形状
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                result = self._find_shape_in_group_layout(shape, text)
                if result:
                    return result
        return None

    def _find_shape_in_group_layout(self, group, text: str):
        """在布局的组合形状中查找"""
        try:
            for shape in group.shapes:
                if hasattr(shape, 'text') and text in shape.text:
                    return shape
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    result = self._find_shape_in_group_layout(shape, text)
                    if result:
                        return result
        except:
            pass
        return None

    def _replace_text_containing_pattern(self, slide, pattern: str, replacement: str) -> bool:
        """
        替换包含特定模式的文本

        Args:
            slide: 幻灯片
            pattern: 要匹配的模式
            replacement: 替换文本
        """
        shape = self._find_shape_containing(slide.shapes, pattern)
        if shape and hasattr(shape, 'text_frame'):
            original_text = shape.text
            if pattern in original_text:
                new_text = original_text.replace(pattern, replacement)
                return self._replace_text_keep_format(shape, original_text, new_text)
        return False

    def _fill_knowledge_content(self, slide, content: str, highlights: List, difficulties: List):
        """
        填充课堂知识内容，带重点/难点颜色标记
        支持多段文本（用换行分隔），每段自动添加数字编号
        同时自动调整白色方框高度以适应内容

        Args:
            slide: 幻灯片
            content: 知识内容文本（多段用换行分隔）
            highlights: 重点词汇列表，支持两种格式：
                        - ["词汇"] - 标记第一次出现
                        - [("词汇", 索引)] - 标记第N次出现
            difficulties: 难点词汇列表，格式同上
        """
        if not content:
            return

        # 查找知识内容区域
        content_shape = self._find_shape_containing(slide.shapes, '理解什么是三维')

        if not content_shape:
            title_shape = self._find_shape_containing(slide.shapes, '课堂知识内容')
            if title_shape and hasattr(title_shape, 'top'):
                content_shape = self._find_shape_below(slide.shapes, title_shape.top, exclude_title=title_shape)

        if not content_shape or not hasattr(content_shape, 'text_frame'):
            return

        tf = content_shape.text_frame
        original_font = self._get_first_run_font(tf)

        # 设置边距和居中对齐
        margin = Emu(int(0.2 * 914400 / 2.54))  # 0.2cm
        bodyPr = tf._bodyPr
        bodyPr.set('tIns', str(int(margin)))
        bodyPr.set('bIns', str(int(margin)))
        bodyPr.set('anchor', 'ctr')  # 上下居中

        # 清除原有段落
        p_list = tf._txBody.findall(qn('a:p'))
        for p in p_list:
            tf._txBody.remove(p)

        # 先用_split_text_by_markers处理整个文本，获取带颜色的分段
        all_segments = self._split_text_by_markers(content, highlights, difficulties)

        # 将分段按行（换行符）组织到段落中
        paragraphs_data = []  # [(段落文本, [(segment文本, 颜色), ...])]
        current_para_segments = []
        current_para_text = ""

        for seg_text, seg_color in all_segments:
            if '\n' in seg_text:
                # 分割换行
                lines = seg_text.split('\n')
                for i, line in enumerate(lines):
                    if i > 0:
                        # 保存当前段落
                        if current_para_segments or current_para_text.strip():
                            paragraphs_data.append((current_para_text.strip(), current_para_segments))
                        current_para_segments = []
                        current_para_text = ""
                    if line:
                        current_para_segments.append((line, seg_color))
                        current_para_text += line
            else:
                current_para_segments.append((seg_text, seg_color))
                current_para_text += seg_text

        # 保存最后一个段落
        if current_para_segments or current_para_text.strip():
            paragraphs_data.append((current_para_text.strip(), current_para_segments))

        # 过滤空段落
        paragraphs_data = [(text, segs) for text, segs in paragraphs_data if text]

        # 添加每段文本（带自动编号）
        for i, (para_text, para_segments) in enumerate(paragraphs_data, 1):
            etree.SubElement(tf._txBody, qn('a:p'))
            para = tf.paragraphs[i - 1]

            # 设置段落属性（编号列表）
            pPr = para._p.get_or_add_pPr()
            pPr.set('marL', '228600')  # 左边距
            pPr.set('indent', '-228600')  # 悬挂缩进

            # 添加自动编号
            buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum'))
            buAutoNum.set('type', 'arabicPeriod')  # 1. 2. 3. ...
            etree.SubElement(pPr, qn('a:buFontTx'))

            # 添加带颜色标记的文本（使用已计算好的分段）
            self._set_colored_text_from_segments(para, para_segments, original_font)

        # 调整高度（同时更新组合的边界框）
        self._adjust_shape_height_via_com(content_shape, margin)

    def _set_colored_text_from_segments(self, paragraph, segments: List[tuple], font_info: dict = None):
        """
        根据预先计算好的分段设置带颜色的文本

        Args:
            paragraph: 段落对象
            segments: [(文本, 颜色), ...] 列表
            font_info: 原始字体信息
        """
        from pptx.util import Pt

        for text, color in segments:
            if not text:
                continue

            run = paragraph.add_run()
            run.text = text

            # 设置颜色 - 默认使用 #404040
            if color:
                run.font.color.rgb = color
            else:
                run.font.color.rgb = DEFAULT_TEXT_COLOR

            # 设置字号
            if font_info and font_info.get('size'):
                run.font.size = font_info['size']
            else:
                run.font.size = Pt(11)

            # 设置字体格式
            if font_info:
                self._set_font_properties(run, font_info)
            else:
                self._set_font_properties(run, {'lang': 'zh-CN', 'alt_lang': 'en-US'})

    def _update_group_bounds(self, child_shape, original_height_emu: int, new_height_emu: int):
        """
        更新包含子形状的组合的边界框

        Args:
            child_shape: 子形状对象
            original_height_emu: 原始高度（EMU）
            new_height_emu: 新的高度（EMU）
        """
        # 计算高度变化量
        height_diff = new_height_emu - original_height_emu

        if height_diff == 0:
            return

        # 获取子形状的 XML 元素
        child_element = child_shape._element

        # 向上查找包含此元素的 grpSp（组合形状）
        current = child_element
        while current is not None:
            parent = current.getparent()
            if parent is None:
                break

            # 检查父元素是否是 grpSp（组合形状）
            if parent.tag == qn('p:grpSp'):
                grpSpPr = parent.find(qn('p:grpSpPr'))
                if grpSpPr is not None:
                    xfrm = grpSpPr.find(qn('a:xfrm'))
                    if xfrm is not None:
                        # 更新组合的视觉大小 (ext)
                        ext = xfrm.find(qn('a:ext'))
                        if ext is not None:
                            old_grp_cy = int(ext.get('cy', 0))
                            new_grp_cy = old_grp_cy + height_diff
                            ext.set('cy', str(new_grp_cy))

                        # 同时更新子形状坐标系的边界框 (chExt)
                        # 这是PowerPoint用来确定子形状总范围的关键属性
                        chExt = xfrm.find(qn('a:chExt'))
                        if chExt is not None:
                            old_ch_cy = int(chExt.get('cy', 0))
                            new_ch_cy = old_ch_cy + height_diff
                            chExt.set('cy', str(new_ch_cy))
                break

            current = parent

    def _set_text_margins(self, shape, top: Emu = None, bottom: Emu = None, left: Emu = None, right: Emu = None):
        """
        设置形状内文本的边距

        Args:
            shape: 形状对象
            top: 上边距（EMU）
            bottom: 下边距（EMU）
            left: 左边距（EMU）
            right: 右边距（EMU）
        """
        if not hasattr(shape, 'text_frame'):
            return

        # 获取 bodyPr 元素
        tf = shape.text_frame
        bodyPr = tf._bodyPr

        if bodyPr is None:
            return

        # 设置边距（单位：EMU）
        if top is not None:
            bodyPr.set('tIns', str(int(top)))
        if bottom is not None:
            bodyPr.set('bIns', str(int(bottom)))
        if left is not None:
            bodyPr.set('lIns', str(int(left)))
        if right is not None:
            bodyPr.set('rIns', str(int(right)))

    def _adjust_shape_height_via_com(self, shape, margin: Emu):
        """
        调整形状高度以适应文本内容

        注意：此方法设置一个足够大的高度，实际高度由COM API后处理调整。
        由于python-pptx无法获取PowerPoint实际渲染的文本高度，
        这里设置一个基于行数估算的初始高度，真正的调整需要COM API。

        Args:
            shape: 形状对象（python-pptx中的矩形）
            margin: 边距（EMU）
        """
        # 保存原始高度（用于计算组合高度变化）
        original_height_emu = shape.height

        # 使用更保守的估算公式（留出足够空间）
        # 基于实际测试数据调整参数
        FONT_SIZE_CM = 11 * 2.54 / 72  # 0.3881 cm
        LINE_SPACING_CM = 0.15  # 调整为0.15cm（之前是0.045cm太小）
        MARGIN_CM = 0.2

        if not hasattr(shape, 'text_frame'):
            return

        tf = shape.text_frame

        # 计算段落数（每个段落一行，因为已经用换行分割）
        total_lines = 0
        for para in tf.paragraphs:
            para_text = ""
            for run in para.runs:
                para_text += run.text
            if para_text.strip():
                total_lines += 1

        total_lines = max(1, total_lines)

        # 计算高度（使用调整后的参数）
        text_height_cm = FONT_SIZE_CM * total_lines + LINE_SPACING_CM * (total_lines - 1)
        rect_height_cm = text_height_cm + 2 * MARGIN_CM
        rect_height_emu = int(rect_height_cm * 914400 / 2.54)

        # 设置新高度
        shape.height = rect_height_emu

        # 更新包含此形状的组合的边界框
        self._update_group_bounds(shape, original_height_emu, rect_height_emu)

    def _calculate_shape_height(self, shape, font_size_cm: float, line_spacing_cm: float, margin_cm: float) -> int:
        """
        计算形状需要的高度（不修改形状）

        已废弃：此方法使用不准确的公式，请使用COM API获取实际高度。

        Args:
            shape: 形状对象
            font_size_cm: 字体大小（cm）
            line_spacing_cm: 行间距（cm）
            margin_cm: 边距（cm）

        Returns:
            计算出的高度（EMU）
        """
        if not hasattr(shape, 'text_frame'):
            return shape.height

        tf = shape.text_frame

        # 基于实际段落数计算
        total_lines = 0
        shape_width_chars = 29  # 每行约29个中文字符

        import math
        for para in tf.paragraphs:
            para_text = ""
            for run in para.runs:
                para_text += run.text

            if not para_text.strip():
                continue

            para_lines = max(1, math.ceil(len(para_text) / shape_width_chars))
            total_lines += para_lines

        total_lines = max(1, total_lines)

        # 计算高度
        text_height_cm = font_size_cm * total_lines + line_spacing_cm * (total_lines - 1)
        rect_height_cm = text_height_cm + 2 * margin_cm
        rect_height_emu = int(rect_height_cm * 914400 / 2.54)

        return rect_height_emu

    def _adjust_shape_height_to_text(self, shape, top_margin: Emu = None, bottom_margin: Emu = None):
        """
        调整形状高度以适应文本内容

        Args:
            shape: 形状对象
            top_margin: 上边距（EMU）
            bottom_margin: 下边距（EMU）
        """
        if not hasattr(shape, 'text_frame'):
            return

        try:
            import win32com.client

            # 使用 COM API 获取实际文本高度
            ppt_app = win32com.client.Dispatch("PowerPoint.Application")
            presentation = ppt_app.ActivePresentation

            if presentation is None:
                return

            # 找到对应的形状
            slide_index = None
            shape_index = None

            for i, sld in enumerate(presentation.Slides):
                for j, shp in enumerate(sld.Shapes):
                    if hasattr(shp, 'TextFrame'):
                        if shp.Name == shape.name:
                            slide_index = i + 1
                            shape_index = j + 1
                            break
                if slide_index:
                    break

            if slide_index and shape_index:
                com_shape = presentation.Slides(slide_index).Shapes(shape_index)

                # 获取文本框的实际边界
                # 使用 TextFrame.TextRange.BoundHeight 获取文本实际高度
                text_height = com_shape.TextFrame.TextRange.BoundHeight

                # 计算新的形状高度
                top_m = int(top_margin) if top_margin else 0
                bottom_m = int(bottom_margin) if bottom_margin else 0

                # 新高度 = 文本高度 + 上边距 + 下边距
                new_height = text_height + top_m + bottom_m

                # 设置新的高度
                shape.height = int(new_height)

        except Exception:
            # 如果 COM 失败，使用估算方法
            self._adjust_shape_height_estimated(shape, top_margin, bottom_margin)

    def _adjust_shape_height_estimated(self, shape, top_margin: Emu = None, bottom_margin: Emu = None):
        """
        使用估算方法调整形状高度（当 COM API 不可用时的备选方案）

        Args:
            shape: 形状对象
            top_margin: 上边距（EMU）
            bottom_margin: 下边距（EMU）
        """
        if not hasattr(shape, 'text_frame'):
            return

        tf = shape.text_frame

        # 获取字体大小
        font_size = None
        if tf.paragraphs and tf.paragraphs[0].runs:
            font_size = tf.paragraphs[0].runs[0].font.size

        if font_size is None:
            font_size = Emu(1100 * 127)  # 默认11pt

        # 估算行高（约1.2倍字体大小）
        line_height = int(font_size * 1.2)

        # 计算文本行数
        total_text = ""
        for para in tf.paragraphs:
            for run in para.runs:
                total_text += run.text

        # 估算每行字符数（根据宽度）
        shape_width_chars = int(shape.width / (font_size * 0.5))  # 中文字符约等于字体大小
        if shape_width_chars < 1:
            shape_width_chars = 30

        # 计算行数
        import math
        num_lines = max(1, math.ceil(len(total_text) / shape_width_chars))

        # 计算新的高度
        top_m = int(top_margin) if top_margin else 0
        bottom_m = int(bottom_margin) if bottom_margin else 0
        new_height = num_lines * line_height + top_m + bottom_m

        # 设置最小高度
        min_height = font_size + top_m + bottom_m
        new_height = max(new_height, min_height)

        shape.height = int(new_height)

    def _get_first_run_font(self, text_frame) -> dict:
        """获取第一个run的字体信息，包括东亚字体"""
        font_info = {
            'name': TEMPLATE_DEFAULT_FONT,
            'name_east': TEMPLATE_DEFAULT_FONT_EAST
        }
        if text_frame.paragraphs:
            para = text_frame.paragraphs[0]
            if para.runs:
                run = para.runs[0]
                if run.font:
                    # 只有当模板中有明确的字体名称时才覆盖
                    if run.font.name:
                        font_info['name'] = run.font.name
                    if run.font.size:
                        font_info['size'] = run.font.size
                    if run.font.bold is not None:
                        font_info['bold'] = run.font.bold
                    if run.font.italic is not None:
                        font_info['italic'] = run.font.italic

                # 从XML获取东亚字体
                rPr = run._r.find(qn('a:rPr'))
                if rPr is not None:
                    cs_elem = rPr.find(qn('a:cs'))
                    if cs_elem is not None:
                        typeface = cs_elem.get('typeface')
                        if typeface:
                            font_info['name_east'] = typeface
        return font_info

    def _find_shape_below(self, shapes, min_top, exclude_title=None):
        """查找位于指定位置下方的形状（用于查找内容区域）"""
        candidates = []

        def search_shapes(shapes_list):
            for shape in shapes_list:
                if hasattr(shape, 'text') and shape.text and hasattr(shape, 'top'):
                    if shape != exclude_title and shape.top > min_top:
                        # 排除评价相关的形状
                        text = shape.text
                        if '优' not in text or '良' not in text:
                            if '课堂表现' not in text:
                                candidates.append((shape.top, shape))
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        search_shapes(shape.shapes)
                    except:
                        pass

        search_shapes(shapes)
        if candidates:
            # 返回最接近的一个
            candidates.sort(key=lambda x: x[0])
            return candidates[0][1]
        return None

    def _set_east_asian_font(self, run, font_name: str):
        """
        设置东亚字体（中文字体）通过XML

        Args:
            run: python-pptx的Run对象
            font_name: 东亚字体名称
        """
        # 获取run的rPr元素
        rPr = run._r.get_or_add_rPr()

        # 设置语言属性（关键！）
        rPr.set('lang', 'zh-CN')
        rPr.set('altLang', 'en-US')

        # 设置东亚字体元素 <a:ea typeface="字体名"/>
        ea_elem = rPr.find(qn('a:ea'))
        if ea_elem is None:
            ea_elem = etree.SubElement(rPr, qn('a:ea'))

        # 设置字体
        ea_elem.set('typeface', font_name)

    def _set_colored_text(self, paragraph, content: str, highlights: List[str], difficulties: List[str], font_info: dict = None):
        """
        在段落中设置带颜色的文本

        Args:
            paragraph: 段落对象
            content: 文本内容
            highlights: 重点词汇列表
            difficulties: 难点词汇列表
            font_info: 原始字体信息（用于保持字号）
        """
        from pptx.util import Pt

        # 创建文本片段列表，每个片段包含文本和颜色
        segments = self._split_text_by_markers(content, highlights, difficulties)

        for text, color in segments:
            run = paragraph.add_run()
            run.text = text

            # 设置颜色 - 默认使用 #404040
            if color:
                run.font.color.rgb = color
            else:
                run.font.color.rgb = DEFAULT_TEXT_COLOR

            # 设置字号 - 优先使用font_info中的字号，否则默认11pt
            if font_info and font_info.get('size'):
                run.font.size = font_info['size']
            else:
                run.font.size = Pt(11)  # 默认11pt

            # 设置字体格式
            if font_info:
                # 设置语言属性（关键）
                self._set_font_properties(run, font_info)
            else:
                # 使用默认语言属性
                self._set_font_properties(run, {'lang': 'zh-CN', 'alt_lang': 'en-US'})

    def _set_font_properties(self, run, font_info: dict):
        """
        设置字体属性，包括语言属性
        注意：只设置latin字体，不设置ea字体，让东亚字体从主题继承

        Args:
            run: python-pptx的Run对象
            font_info: 字体信息字典
        """
        rPr = run._r.get_or_add_rPr()

        # 设置语言属性（关键！）
        rPr.set('lang', font_info.get('lang', 'zh-CN'))
        rPr.set('altLang', font_info.get('alt_lang', 'en-US'))

        # 只设置拉丁字体，东亚字体从主题继承（和模板一致）
        latin_elem = rPr.find(qn('a:latin'))
        if latin_elem is None:
            latin_elem = etree.SubElement(rPr, qn('a:latin'))
        latin_elem.set('typeface', font_info.get('name', TEMPLATE_DEFAULT_FONT))

    def _split_text_by_markers(self, content: str, highlights: List, difficulties: List) -> List[tuple]:
        """
        根据重点/难点标记分割文本
        支持按出现索引标记特定位置的词汇

        Args:
            content: 原始文本
            highlights: 重点词汇列表，支持两种格式：
                        - ["词汇"] - 标记第一次出现
                        - [("词汇", 索引)] - 标记第N次出现（索引从1开始）
            difficulties: 难点词汇列表，格式同上

        Returns:
            [(文本, 颜色), ...] 列表
        """
        if not highlights and not difficulties:
            return [(content, None)]

        def find_nth_occurrence(text: str, word: str, n: int) -> int:
            """查找词汇的第N次出现位置（n从1开始）"""
            start = 0
            for i in range(n):
                pos = text.find(word, start)
                if pos == -1:
                    return -1
                if i == n - 1:  # 找到第n次时返回
                    return pos
                start = pos + len(word)  # 跳过已找到的词汇
            return -1

        def parse_marker(marker):
            """解析标记，返回(词汇, 出现索引)"""
            if isinstance(marker, tuple):
                return marker[0], marker[1]  # (词汇, 索引)
            else:
                return marker, 1  # 默认标记第一次出现

        # 创建标记位置列表
        markers = []
        for marker in highlights:
            word, occurrence = parse_marker(marker)
            pos = find_nth_occurrence(content, word, occurrence)
            if pos != -1:
                markers.append((pos, pos + len(word), word, HIGHLIGHT_COLOR))

        for marker in difficulties:
            word, occurrence = parse_marker(marker)
            pos = find_nth_occurrence(content, word, occurrence)
            if pos != -1:
                markers.append((pos, pos + len(word), word, DIFFICULTY_COLOR))

        # 按位置排序
        markers.sort(key=lambda x: x[0])

        # 去除重叠的标记（保留先出现的）
        filtered_markers = []
        last_end = 0
        for start, end, word, color in markers:
            if start >= last_end:
                filtered_markers.append((start, end, word, color))
                last_end = end

        # 分割文本
        segments = []
        last_pos = 0

        for start, end, word, color in filtered_markers:
            # 添加前面的普通文本
            if start > last_pos:
                segments.append((content[last_pos:start], None))
            # 添加带颜色的文本
            segments.append((word, color))
            last_pos = end

        # 添加剩余的文本
        if last_pos < len(content):
            segments.append((content[last_pos:], None))

        return segments if segments else [(content, None)]

    def _replace_shape_containing(self, slide, search_text: str, new_text: str, is_label: bool = False) -> bool:
        """
        替换包含特定文本的形状内容

        Args:
            slide: 幻灯片
            search_text: 搜索文本
            new_text: 新文本
            is_label: 是否是标签（如果是标签，则找该形状后面的形状）
        """
        if is_label:
            # 找到标签形状，然后找它附近的值形状
            # 这个逻辑比较复杂，简化处理：直接搜索并替换
            return False

        shape = self._find_shape_containing(slide.shapes, search_text)
        if shape and new_text:
            return self._replace_text_keep_format(shape, shape.text, new_text)
        return False

    def fill_model_images(self, image_paths: List[str], start_slide: int = 2) -> int:
        """
        填充模型展示页图片

        Args:
            image_paths: 图片路径列表
            start_slide: 开始的幻灯片索引

        Returns:
            成功填充的图片数量
        """
        filled = 0
        slide_index = start_slide

        for img_path in image_paths:
            if not Path(img_path).exists():
                continue

            if slide_index >= len(self.prs.slides):
                break

            slide = self.prs.slides[slide_index]

            # 移除现有图片
            self._remove_pictures(slide)

            # 插入新图片
            if self._insert_model_image(slide, img_path):
                filled += 1

            slide_index += 1

        return filled

    def fill_work_images(self, image_paths: List[str], start_slide: int = 3) -> int:
        """
        填充作品展示页图片

        支持多张图片，每张图片对应一页作品展示幻灯片。
        图片按上传顺序填充到连续的作品展示页中。

        Args:
            image_paths: 图片路径列表
            start_slide: 开始的幻灯片索引

        Returns:
            成功填充的图片数量
        """
        filled = 0
        slide_index = start_slide

        for img_path in image_paths:
            if not Path(img_path).exists():
                continue

            if slide_index >= len(self.prs.slides):
                break

            slide = self.prs.slides[slide_index]

            # 移除现有图片
            self._remove_pictures(slide)

            # 插入新图片
            if self._insert_work_image(slide, img_path):
                filled += 1

            slide_index += 1

        return filled

    def fill_program_images(self, image_paths: List[str], start_slide: int) -> int:
        """
        填充程序展示页图片

        Args:
            image_paths: 图片路径列表
            start_slide: 开始的幻灯片索引

        Returns:
            成功填充的图片数量
        """
        filled = 0
        slide_index = start_slide

        for img_path in image_paths:
            if not Path(img_path).exists():
                continue

            if slide_index >= len(self.prs.slides):
                break

            slide = self.prs.slides[slide_index]

            # 移除现有图片
            self._remove_pictures(slide)

            # 插入新图片（使用程序展示专用尺寸：16cm x 20cm）
            if self._insert_program_image(slide, img_path):
                filled += 1

            slide_index += 1

        return filled

    def fill_vehicle_images(self, image_paths: List[str], start_slide: int) -> int:
        """
        填充车辆展示页图片

        Args:
            image_paths: 图片路径列表
            start_slide: 开始的幻灯片索引

        Returns:
            成功填充的图片数量
        """
        filled = 0
        slide_index = start_slide

        for img_path in image_paths:
            if not Path(img_path).exists():
                continue

            if slide_index >= len(self.prs.slides):
                break

            slide = self.prs.slides[slide_index]

            # 移除现有图片
            self._remove_pictures(slide)

            # 插入新图片（使用与work相同的尺寸）
            if self._insert_work_image(slide, img_path):
                filled += 1

            slide_index += 1

        return filled

    def _replace_text_in_slide(self, slide, old_text: str, new_text: str) -> bool:
        """在幻灯片中查找并替换文本"""
        shape = self._find_shape_containing(slide.shapes, old_text)
        if shape:
            return self._replace_text_keep_format(shape, old_text, new_text)
        return False

    def _replace_lesson_number(self, slide, old_text: str, new_text: str) -> bool:
        """替换课时编号（加粗，颜色#D3F1F0）"""
        shape = self._find_shape_containing(slide.shapes, old_text)
        if shape and hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            # 替换文本
            shape.text = shape.text.replace(old_text, new_text)

            # 设置包含"第"和"课"的run的样式
            TARGET_COLOR = RGBColor(0xD3, 0xF1, 0xF0)
            for para in tf.paragraphs:
                for run in para.runs:
                    run_text = run.text
                    if "第" in run_text and "课" in run_text:
                        if run.font:
                            run.font.bold = True
                            run.font.color.rgb = TARGET_COLOR
            return True
        return False

    def _replace_text_containing(self, slide, search_text: str, new_text: str) -> bool:
        """替换包含特定文本的形状的全部内容"""
        shape = self._find_shape_containing(slide.shapes, search_text)
        if shape and new_text:
            return self._replace_text_keep_format(shape, shape.text, new_text)
        return False

    def _find_shape_containing(self, shapes, text: str):
        """递归查找包含指定文本的形状"""
        for shape in shapes:
            # 先检查是否是组合形状
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                result = self._find_shape_in_group(shape, text)
                if result:
                    return result
            # 检查普通形状的文本
            elif hasattr(shape, 'text') and text in shape.text:
                return shape
        return None

    def _find_shape_in_group(self, group, text: str):
        """在组合形状中查找"""
        try:
            for shape in group.shapes:
                if hasattr(shape, 'text') and text in shape.text:
                    return shape
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    result = self._find_shape_in_group(shape, text)
                    if result:
                        return result
        except:
            pass
        return None

    def _replace_text_keep_format(self, shape, old_text: str, new_text: str, use_default_color: bool = True) -> bool:
        """
        替换文本但保留格式

        Args:
            shape: 形状对象
            old_text: 要替换的文本
            new_text: 新文本
            use_default_color: 是否使用默认颜色(#404040)
        """
        if not hasattr(shape, 'text'):
            return False

        if old_text not in shape.text:
            return False

        # 保存格式信息
        font_info = {}
        if hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            if len(tf.paragraphs) > 0 and len(tf.paragraphs[0].runs) > 0:
                first_run = tf.paragraphs[0].runs[0]
                if first_run.font:
                    font_info = {
                        'name': first_run.font.name,
                        'size': first_run.font.size,
                        'bold': first_run.font.bold,
                        'italic': first_run.font.italic,
                    }

        # 替换文本
        if old_text == shape.text:
            shape.text = new_text
        else:
            shape.text = shape.text.replace(old_text, new_text)

        # 恢复格式并设置颜色
        if hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    if run.font:
                        if font_info.get('name'):
                            run.font.name = font_info['name']
                        if font_info.get('size'):
                            run.font.size = font_info['size']
                        if font_info.get('bold') is not None:
                            run.font.bold = font_info['bold']
                        if font_info.get('italic') is not None:
                            run.font.italic = font_info['italic']
                        # 设置默认文字颜色
                        if use_default_color:
                            run.font.color.rgb = DEFAULT_TEXT_COLOR

        return True

    def _remove_pictures(self, slide):
        """移除幻灯片中的所有图片"""
        for shape in list(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                sp = shape._element
                sp.getparent().remove(sp)

    def _insert_model_image(self, slide, image_path: str) -> bool:
        """
        插入模型展示页图片

        要求：
        - 保持原比例，通过裁剪变成4:5比例
        - 目标尺寸：高度20cm，宽度16cm
        - 上下位置使用模板预定义，左右居中
        """
        try:
            from PIL import Image
            from pptx.util import Cm
            import io

            # 目标比例：4:5
            target_ratio = 4 / 5

            # 目标尺寸：高度20cm，宽度16cm
            target_height_cm = 20
            target_width_cm = 16

            # 垂直位置：使用模板预定义的top位置 + 0.4cm
            # 3.62cm + 0.4cm = 4.02cm
            top_pos = Cm(4.00)

            # 加载原始图片
            img = Image.open(image_path)
            orig_width, orig_height = img.size
            orig_ratio = orig_width / orig_height

            # 计算裁剪区域（居中裁剪到4:5比例）
            if orig_ratio > target_ratio:
                # 原图更宽，需要裁剪宽度
                new_width = int(orig_height * target_ratio)
                left = (orig_width - new_width) // 2
                crop_top = 0
                right = left + new_width
                bottom = orig_height
            else:
                # 原图更高，需要裁剪高度
                new_height = int(orig_width / target_ratio)
                left = 0
                crop_top = (orig_height - new_height) // 2
                right = orig_width
                bottom = crop_top + new_height

            # 裁剪图片
            cropped = img.crop((left, crop_top, right, bottom))

            # 保存到内存
            img_buffer = io.BytesIO()
            # 保持原始格式
            img_format = img.format if img.format else 'PNG'
            if img_format == 'JPEG' or image_path.lower().endswith(('.jpg', '.jpeg')):
                cropped.save(img_buffer, format='JPEG', quality=95)
            else:
                cropped.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            # 移除现有图片
            self._remove_pictures(slide)

            # 获取幻灯片尺寸
            slide_width = self.prs.slide_width

            # 计算水平居中位置
            # 幻灯片宽度 - 图片宽度 = 剩余空间，除以2得到左边距
            left_pos = (slide_width - Cm(target_width_cm)) / 2

            # 插入图片：指定宽度和高度
            slide.shapes.add_picture(
                img_buffer,
                left_pos, top_pos,
                width=Cm(target_width_cm),
                height=Cm(target_height_cm)
            )

            return True

        except Exception as e:
            print(f"插入模型图片失败: {e}")
            import traceback
            traceback.print_exc()
            return False

    def _insert_work_image(self, slide, image_path: str) -> bool:
        """
        插入精彩瞬间/车辆展示页图片

        要求：
        - 保持原比例，通过裁剪变成4:5比例
        - 目标尺寸：高度21cm，宽度16.8cm
        - 垂直位置3.60cm，水平居中
        """
        try:
            from PIL import Image
            from pptx.util import Cm
            import io

            # 目标比例：4:5
            target_ratio = 4 / 5

            # 目标尺寸：高度21cm，宽度16.8cm
            target_height_cm = 21
            target_width_cm = 16.8

            # 垂直位置：3.60cm
            top_pos = Cm(3.60)

            # 加载原始图片
            img = Image.open(image_path)
            orig_width, orig_height = img.size
            orig_ratio = orig_width / orig_height

            # 计算裁剪区域（居中裁剪到4:5比例）
            if orig_ratio > target_ratio:
                # 原图更宽，需要裁剪宽度
                new_width = int(orig_height * target_ratio)
                left = (orig_width - new_width) // 2
                crop_top = 0
                right = left + new_width
                bottom = orig_height
            else:
                # 原图更高，需要裁剪高度
                new_height = int(orig_width / target_ratio)
                left = 0
                crop_top = (orig_height - new_height) // 2
                right = orig_width
                bottom = crop_top + new_height

            # 裁剪图片
            cropped = img.crop((left, crop_top, right, bottom))

            # 保存到内存
            img_buffer = io.BytesIO()
            # 保持原始格式
            img_format = img.format if img.format else 'PNG'
            if img_format == 'JPEG' or image_path.lower().endswith(('.jpg', '.jpeg')):
                cropped.save(img_buffer, format='JPEG', quality=95)
            else:
                cropped.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            # 移除现有图片
            self._remove_pictures(slide)

            # 获取幻灯片尺寸
            slide_width = self.prs.slide_width

            # 计算水平居中位置
            left_pos = (slide_width - Cm(target_width_cm)) / 2

            # 插入图片：指定宽度和高度
            slide.shapes.add_picture(
                img_buffer,
                left_pos, top_pos,
                width=Cm(target_width_cm),
                height=Cm(target_height_cm)
            )

            return True

        except Exception as e:
            print(f"插入作品图片失败: {e}")
            import traceback
            traceback.print_exc()
            return False

    def _insert_program_image(self, slide, image_path: str) -> bool:
        """
        插入程序展示页图片

        要求：
        - 保持原始比例，不裁剪
        - 偏高图片：高度=20cm，宽度按比例缩放
        - 偏宽图片：宽度=16cm，高度按比例缩放
        - 垂直位置4.00cm，水平居中
        """
        try:
            from PIL import Image
            from pptx.util import Cm
            import io

            # 目标区域最大尺寸
            max_width_cm = 16
            max_height_cm = 20

            # 4:5 = 目标区域比例
            target_ratio = 4 / 5

            # 垂直位置：4.00cm
            top_pos = Cm(4.00)

            # 加载原始图片
            img = Image.open(image_path)
            orig_width, orig_height = img.size
            orig_ratio = orig_width / orig_height

            # 根据图片比例选择缩放基准
            if orig_ratio < target_ratio:
                # 偏高图片：以高度为基准
                actual_height_cm = max_height_cm
                actual_width_cm = max_height_cm * orig_ratio
            else:
                # 偏宽或正好的图片：以宽度为基准
                actual_width_cm = max_width_cm
                actual_height_cm = max_width_cm / orig_ratio

            # 保存到内存（不裁剪，保持原图）
            img_buffer = io.BytesIO()
            img_format = img.format if img.format else 'PNG'
            if img_format == 'JPEG' or image_path.lower().endswith(('.jpg', '.jpeg')):
                img.save(img_buffer, format='JPEG', quality=95)
            else:
                img.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            # 移除现有图片
            self._remove_pictures(slide)

            # 获取幻灯片尺寸
            slide_width = self.prs.slide_width

            # 计算水平居中位置
            left_pos = (slide_width - Cm(actual_width_cm)) / 2

            # 插入图片：指定实际宽度和高度
            slide.shapes.add_picture(
                img_buffer,
                left_pos, top_pos,
                width=Cm(actual_width_cm),
                height=Cm(actual_height_cm)
            )

            return True

        except Exception as e:
            print(f"插入程序图片失败: {e}")
            import traceback
            traceback.print_exc()
            return False


def fill_ppt_content(presentation: Presentation, data: CourseUnitData,
                      series_name: str = None, series_level: int = None,
                      model_image_count: int = 0, work_image_count: int = 0,
                      program_image_count: int = 0, vehicle_image_count: int = 0,
                      include_cover: bool = True) -> bool:
    """
    填充PPT内容的便捷函数

    页面顺序：
    - 封面 (索引0，仅当include_cover=True)
    - 课程信息 (索引0或1，取决于是否有封面)
    - 模型展示页
    - 程序展示页
    - 车辆展示页
    - 精彩瞬间

    Args:
        presentation: PPT演示文稿对象
        data: 课程数据
        series_name: 课程系列名称
        series_level: 课程阶数
        model_image_count: 模型展示页数量
        work_image_count: 精彩瞬间页数量
        program_image_count: 程序展示页数量
        vehicle_image_count: 车辆展示页数量
        include_cover: 是否包含封面页（第1课为True，其他为False）

    Returns:
        是否成功
    """
    filler = ContentFiller(presentation)

    # 填充封面页的系列名称（仅当包含封面时）
    if include_cover and series_name and series_level:
        filler.fill_cover_series(series_name, series_level)

    # 计算课程信息页索引
    course_info_index = 1 if include_cover else 0

    # 填充课程信息（传递正确的幻灯片索引）
    filler.fill_course_info(data, series_name, series_level, slide_index=course_info_index)

    # 计算各类型页面的起始索引
    # 如果有封面：幻灯片索引 0=封面，1=课程信息，2+=内容页
    # 如果无封面：幻灯片索引 0=课程信息，1+=内容页
    base_index = 2 if include_cover else 1

    # 填充模型图片
    if data.model_images and model_image_count > 0:
        filler.fill_model_images(data.model_images, start_slide=base_index)
    base_index += model_image_count

    # 填充程序展示图片
    if data.program_images and program_image_count > 0:
        filler.fill_program_images(data.program_images, start_slide=base_index)
    base_index += program_image_count

    # 填充车辆展示图片
    if data.vehicle_images and vehicle_image_count > 0:
        filler.fill_vehicle_images(data.vehicle_images, start_slide=base_index)
    base_index += vehicle_image_count

    # 填充精彩瞬间图片
    if data.work_images and work_image_count > 0:
        filler.fill_work_images(data.work_images, start_slide=base_index)

    return True


def _find_course_info_slide(pres) -> int:
    """
    在PPT中查找课程信息页的索引

    课程信息页的特征：包含"学生姓名"或"课程内容"等文本

    Args:
        pres: PowerPoint Presentation 对象

    Returns:
        课程信息页的索引（1-based），如果找不到则返回1
    """
    try:
        for i in range(1, pres.Slides.Count + 1):
            slide = pres.Slides(i)
            for j in range(1, slide.Shapes.Count + 1):
                shape = slide.Shapes(j)
                if shape.HasTextFrame:
                    text = shape.TextFrame.TextRange.Text
                    # 课程信息页的特征文本
                    if any(keyword in text for keyword in ["学生姓名", "授课教师", "课程内容", "上课时间"]):
                        return i
        # 默认返回第1页
        return 1
    except:
        return 1


def post_process_ppt(file_path: str, course_info_slide_index: int = None, keep_open: bool = True) -> bool:
    """
    一次性完成所有 PPT 后处理（只打开/关闭 PPT 一次）

    包含的处理：
    1. 课堂知识内容高度自适应
    2. 课堂评价补充说明高度自适应
    3. 组合纵向分布

    Args:
        file_path: PPT文件路径
        course_info_slide_index: 课程信息页索引（1-based），如果为None则自动查找
        keep_open: 处理完成后是否保持 PPT 打开状态（默认 True）

    Returns:
        是否成功
    """
    MARGIN_CM = 0.2  # 知识内容边距
    MARGIN_TOP_CM = 0.16  # 补充说明上边距
    MARGIN_BOTTOM_CM = 0.15  # 补充说明下边距

    try:
        import win32com.client
        import os

        file_path = os.path.abspath(file_path)

        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True

        pres = ppt_app.Presentations.Open(file_path)

        # 自动查找课程信息页
        if course_info_slide_index is None:
            course_info_slide_index = _find_course_info_slide(pres)

        slide = pres.Slides(course_info_slide_index)

        # ========== 1. 调整知识内容高度 ==========
        for shape_idx in range(1, slide.Shapes.Count + 1):
            grp = slide.Shapes(shape_idx)
            if grp.Type == 6:  # msoGroup
                is_knowledge_group = False
                rect_item = None

                for item_idx in range(1, grp.GroupItems.Count + 1):
                    item = grp.GroupItems.Item(item_idx)
                    if item.HasTextFrame:
                        text = item.TextFrame.TextRange.Text
                        if "课堂知识内容" in text:
                            is_knowledge_group = True
                        elif "Rectangle" in item.Name and len(text) > 20:
                            rect_item = item

                if is_knowledge_group and rect_item:
                    text_height_pt = rect_item.TextFrame.TextRange.BoundHeight
                    margin_pt = MARGIN_CM / 2.54 * 72
                    new_height_pt = text_height_pt + 2 * margin_pt
                    rect_item.TextFrame.MarginTop = margin_pt
                    rect_item.TextFrame.MarginBottom = margin_pt
                    rect_item.Height = new_height_pt
                    print(f"[后处理] 知识内容高度调整: 文本={text_height_pt:.1f}pt, 总高={new_height_pt:.1f}pt")
                    break

        # ========== 2. 调整补充说明高度 ==========
        for shape_idx in range(1, slide.Shapes.Count + 1):
            grp = slide.Shapes(shape_idx)
            if grp.Type == 6:  # msoGroup
                has_additional_comments = False
                for item_idx in range(1, grp.GroupItems.Count + 1):
                    item = grp.GroupItems.Item(item_idx)
                    if item.HasTextFrame:
                        text = item.TextFrame.TextRange.Text
                        if "补充说明" in text:
                            has_additional_comments = True
                            break

                if has_additional_comments:
                    text_box = None
                    gray_frame = None

                    for item_idx in range(1, grp.GroupItems.Count + 1):
                        item = grp.GroupItems.Item(item_idx)
                        item_name = item.Name

                        if item.HasTextFrame:
                            text = item.TextFrame.TextRange.Text
                            is_evaluation = "优" in text and "良" in text
                            is_label = "补充说明" in text and len(text) < 10
                            if len(text) > 10 and not is_evaluation and not is_label:
                                text_box = item

                        if "矩形" in item_name or "Rectangle" in item_name:
                            if not item.HasTextFrame or len(item.TextFrame.TextRange.Text) < 10:
                                width_cm = item.Width / 72 * 2.54
                                if width_cm > 15:
                                    gray_frame = item

                    if text_box and gray_frame:
                        white_frame = None
                        for item_idx in range(1, grp.GroupItems.Count + 1):
                            item = grp.GroupItems.Item(item_idx)
                            item_name = item.Name
                            if "Rectangle" in item_name or "矩形" in item_name:
                                if item_name != gray_frame.Name:
                                    width_cm = item.Width / 72 * 2.54
                                    if width_cm > 15:
                                        white_frame = item
                                        break

                        gray_frame_top = gray_frame.Top
                        text_height_pt = text_box.TextFrame.TextRange.BoundHeight
                        margin_top_pt = MARGIN_TOP_CM / 2.54 * 72
                        margin_bottom_pt = MARGIN_BOTTOM_CM / 2.54 * 72
                        new_text_height_pt = text_height_pt + margin_top_pt + margin_bottom_pt

                        text_box.TextFrame.MarginTop = margin_top_pt
                        text_box.TextFrame.MarginBottom = margin_bottom_pt
                        text_box.Height = new_text_height_pt

                        text_box_bottom = text_box.Top + text_box.Height
                        new_gray_height = text_box_bottom - gray_frame_top
                        gray_frame.Height = new_gray_height

                        if white_frame:
                            white_frame_top = white_frame.Top
                            new_white_height = gray_frame_top - white_frame_top
                            if new_white_height > 0:
                                white_frame.Height = new_white_height

                        print(f"[后处理] 补充说明高度调整: 文本={text_height_pt:.1f}pt, 灰色框={new_gray_height/72*2.54:.2f}cm")
                        break

        # ========== 3. 纵向分布组合 ==========
        groups = {}
        for i in range(1, slide.Shapes.Count + 1):
            shape = slide.Shapes(i)
            if shape.Type == 6:  # msoGroup
                try:
                    for j in range(1, shape.GroupItems.Count + 1):
                        item = shape.GroupItems.Item(j)
                        if item.HasTextFrame:
                            text = item.TextFrame.TextRange.Text
                            if "学生姓名" in text:
                                groups["top"] = shape
                                break
                            elif "课时" in text and "上课时间" not in text:
                                groups["time"] = shape
                                break
                            elif "课堂知识内容" in text:
                                groups["knowledge"] = shape
                                break
                            elif "课堂表现" in text:
                                groups["evaluation"] = shape
                                break
                            elif "课程作业" in text:
                                groups["bottom"] = shape
                                break
                except:
                    pass

        required_keys = ["top", "time", "knowledge", "evaluation", "bottom"]
        if all(key in groups for key in required_keys):
            top_anchor = groups["top"]
            bottom_anchor = groups["bottom"]

            top_bottom = top_anchor.Top + top_anchor.Height
            bottom_top = bottom_anchor.Top

            middle_groups = [groups["time"], groups["knowledge"], groups["evaluation"]]
            total_middle_height = sum(g.Height for g in middle_groups)
            available_space = bottom_top - top_bottom
            gap = (available_space - total_middle_height) / 4

            current_top = top_bottom + gap
            for group in middle_groups:
                group.Top = current_top
                current_top += group.Height + gap

            print(f"[后处理] 纵向分布完成: 间距={gap/72*2.54:.2f}cm")
        else:
            missing = [k for k in required_keys if k not in groups]
            print(f"[后处理] 警告: 未找到组合 {missing}，跳过纵向分布")

        # 保存
        pres.Save()

        # ========== 4. 导出PDF ==========
        try:
            from src.core.ppt_generator import PPTGenerator

            # 生成PDF文件名（去除时间戳）
            ppt_filename = os.path.basename(file_path)
            pdf_filename = PPTGenerator.generate_pdf_filename(ppt_filename)
            pdf_path = os.path.join(os.path.dirname(file_path), pdf_filename)

            # ExportAsFixedFormat 参数说明
            # ppFixedFormatTypePDF = 2 (PDF格式)
            # ppFixedFormatIntentScreen = 1, ppFixedFormatIntentPrint = 2 (打印质量/高分辨率)
            # PrintRange=None 是必须的 - 这是 PowerPoint COM API 的 bug workaround
            # 参考: https://stackoverflow.com/questions/17896216
            pres.ExportAsFixedFormat(
                pdf_path,
                FixedFormatType=2,   # ppFixedFormatTypePDF
                Intent=2,            # ppFixedFormatIntentPrint (高分辨率)
                PrintRange=None      # 必须！解决 "cannot be converted to a COM object" 错误
            )
            print(f"[后处理] PDF已导出: {pdf_path}")

        except Exception as pdf_error:
            # PDF导出失败不影响主流程
            print(f"[后处理] 警告: PDF导出失败 - {pdf_error}")

        if keep_open:
            print("[后处理] PPT 已保存并保持打开状态")
        else:
            pres.Close()
            print("[后处理] PPT 已保存并关闭")

        return True

    except Exception as e:
        print(f"[后处理] 失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def adjust_shape_height_by_text(file_path: str, course_info_slide_index: int = None) -> bool:
    """
    根据实际文本高度调整知识内容区域的形状高度（COM API后处理）

    此函数必须在PPT保存后调用，使用COM API获取实际渲染的文本高度，
    然后设置形状高度 = 文本高度 + 2 * 边距(0.4cm)

    Args:
        file_path: PPT文件路径
        course_info_slide_index: 课程信息页索引（1-based），如果为None则自动查找

    Returns:
        是否成功
    """
    MARGIN_CM = 0.2  # 边距0.2cm

    try:
        import win32com.client
        import os

        file_path = os.path.abspath(file_path)

        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True

        pres = ppt_app.Presentations.Open(file_path)

        # 自动查找课程信息页
        if course_info_slide_index is None:
            course_info_slide_index = _find_course_info_slide(pres)

        slide = pres.Slides(course_info_slide_index)

        for shape_idx in range(1, slide.Shapes.Count + 1):
            grp = slide.Shapes(shape_idx)
            if grp.Type == 6:  # msoGroup
                # 检查是否是知识内容组合
                is_knowledge_group = False
                rect_item = None

                for item_idx in range(1, grp.GroupItems.Count + 1):
                    item = grp.GroupItems.Item(item_idx)
                    if item.HasTextFrame:
                        text = item.TextFrame.TextRange.Text
                        if "课堂知识内容" in text:
                            is_knowledge_group = True
                        elif "Rectangle" in item.Name and len(text) > 20:
                            # 这是内容矩形（文本较长）
                            rect_item = item

                if is_knowledge_group and rect_item:
                    # 获取实际文本高度
                    text_height_pt = rect_item.TextFrame.TextRange.BoundHeight

                    # 计算新的形状高度 = 文本高度 + 2 * 边距
                    margin_pt = MARGIN_CM / 2.54 * 72  # 0.2cm → points
                    new_height_pt = text_height_pt + 2 * margin_pt

                    # 设置边距（确保是0.2cm）
                    rect_item.TextFrame.MarginTop = margin_pt
                    rect_item.TextFrame.MarginBottom = margin_pt

                    # 设置高度
                    rect_item.Height = new_height_pt

                    print(f"调整高度: 文本={text_height_pt:.1f}pt, 边距={margin_pt:.1f}pt, 总高={new_height_pt:.1f}pt")
                    break

        pres.Save()
        pres.Close()

        return True

    except Exception as e:
        print(f"调整高度失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def distribute_groups_vertically(file_path: str, course_info_slide_index: int = None) -> bool:
    """
    纵向均匀分布课程信息页中的组合框（后处理函数）
    需要在PPT文件保存后调用，使用COM API

    分布规则：
    - 组合294（学生姓名、授课教师、课程内容）：位置不变（顶部锚点）
    - 组合3（课程作业、注意事项）：位置不变（底部锚点）
    - 中间三个组合（组合295、组合7、组合5）在顶部和底部之间均匀分布

    Args:
        file_path: PPT文件路径
        course_info_slide_index: 课程信息页索引（1-based），如果为None则自动查找

    Returns:
        是否成功
    """
    try:
        import win32com.client
        import os

        file_path = os.path.abspath(file_path)

        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True

        # 打开文件
        pres = ppt_app.Presentations.Open(file_path)

        # 自动查找课程信息页
        if course_info_slide_index is None:
            course_info_slide_index = _find_course_info_slide(pres)

        slide = pres.Slides(course_info_slide_index)

        # 找到5个组合框
        groups = {}
        for i in range(1, slide.Shapes.Count + 1):
            shape = slide.Shapes(i)
            if shape.Type == 6:  # msoGroup
                # 检查组合内容来识别
                try:
                    for j in range(1, shape.GroupItems.Count + 1):
                        item = shape.GroupItems.Item(j)
                        if item.HasTextFrame:
                            text = item.TextFrame.TextRange.Text
                            if "学生姓名" in text:
                                groups["top"] = shape  # 组合294 - 顶部锚点
                                break
                            elif "课时" in text and "上课时间" not in text:
                                groups["time"] = shape  # 组合295 - 时间信息
                                break
                            elif "课堂知识内容" in text:
                                groups["knowledge"] = shape  # 组合7 - 课堂知识
                                break
                            elif "课堂表现" in text:
                                groups["evaluation"] = shape  # 组合5 - 课堂表现
                                break
                            elif "课程作业" in text:
                                groups["bottom"] = shape  # 组合3 - 底部锚点
                                break
                except:
                    pass

        # 验证找到了所有组合
        required_keys = ["top", "time", "knowledge", "evaluation", "bottom"]
        for key in required_keys:
            if key not in groups:
                print(f"警告: 未找到 {key} 组合")

        if len(groups) < 5:
            print(f"错误: 只找到 {len(groups)} 个组合，需要5个")
            pres.Close()
            return False

        # 获取顶部和底部锚点位置
        top_anchor = groups["top"]
        bottom_anchor = groups["bottom"]

        top_bottom = top_anchor.Top + top_anchor.Height  # 顶部组合的底部
        bottom_top = bottom_anchor.Top  # 底部组合的顶部

        # 计算中间三个组合的总高度
        middle_groups = [groups["time"], groups["knowledge"], groups["evaluation"]]
        total_middle_height = sum(g.Height for g in middle_groups)

        # 计算可用空间
        available_space = bottom_top - top_bottom

        # 计算间距（4个间隙：top到time, time到knowledge, knowledge到evaluation, evaluation到bottom）
        gap = (available_space - total_middle_height) / 4

        # 设置新位置
        current_top = top_bottom + gap

        for group in middle_groups:
            group.Top = current_top
            current_top += group.Height + gap

        # 保存并关闭
        pres.Save()
        pres.Close()

        print("纵向分布完成")
        return True

    except Exception as e:
        print(f"纵向分布失败: {e}")
        import traceback
        traceback.print_exc()
        return False


def adjust_additional_comments_height(file_path: str, course_info_slide_index: int = None) -> bool:
    """
    调整补充说明区域的高度（COM API后处理）

    调整规则：
    - 灰色框上边沿位置不动
    - 文本框上边沿位置不动
    - 文本框高度根据内部文本实际高度自动调整
    - 灰色框底部与文本框底部保持对齐

    Args:
        file_path: PPT文件路径
        course_info_slide_index: 课程信息页索引（1-based），如果为None则自动查找

    Returns:
        是否成功
    """
    MARGIN_CM = 0.16  # 上边距保持原来的0.16cm
    MARGIN_BOTTOM_CM = 0.15  # 下边距0.15cm

    try:
        import win32com.client
        import os

        file_path = os.path.abspath(file_path)

        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = True

        pres = ppt_app.Presentations.Open(file_path)

        # 自动查找课程信息页
        if course_info_slide_index is None:
            course_info_slide_index = _find_course_info_slide(pres)

        slide = pres.Slides(course_info_slide_index)

        # 找到课堂表现评价区域（组合5）
        for shape_idx in range(1, slide.Shapes.Count + 1):
            grp = slide.Shapes(shape_idx)
            if grp.Type == 6:  # msoGroup
                # 检查是否包含"补充说明"
                has_additional_comments = False
                for item_idx in range(1, grp.GroupItems.Count + 1):
                    item = grp.GroupItems.Item(item_idx)
                    if item.HasTextFrame:
                        text = item.TextFrame.TextRange.Text
                        if "补充说明" in text:
                            has_additional_comments = True
                            break

                if has_additional_comments:
                    # 找到文本框和灰色框
                    text_box = None
                    gray_frame = None

                    for item_idx in range(1, grp.GroupItems.Count + 1):
                        item = grp.GroupItems.Item(item_idx)
                        item_name = item.Name

                        # 找包含补充说明内容的文本框
                        # 条件：有文本、文本长度>10、不是评价项、不是"补充说明"标签
                        if item.HasTextFrame:
                            text = item.TextFrame.TextRange.Text
                            # 排除评价项（包含"优"和"良"的）
                            is_evaluation = "优" in text and "良" in text
                            # 排除"补充说明"标签
                            is_label = "补充说明" in text and len(text) < 10
                            # 文本长度>10且不是评价项或标签
                            if len(text) > 10 and not is_evaluation and not is_label:
                                text_box = item

                        # 找灰色背景框（矩形，无文本或文本很短）
                        if "矩形" in item_name or "Rectangle" in item_name:
                            if not item.HasTextFrame or len(item.TextFrame.TextRange.Text) < 10:
                                # 检查尺寸，灰色框应该是16.82cm宽
                                width_cm = item.Width / 72 * 2.54  # points to cm
                                if width_cm > 15:  # 大于15cm的矩形
                                    gray_frame = item

                    if text_box and gray_frame:
                        # 找到白色背景框（Rectangle 44，宽度>15cm且不是灰色框）
                        white_frame = None
                        for item_idx in range(1, grp.GroupItems.Count + 1):
                            item = grp.GroupItems.Item(item_idx)
                            item_name = item.Name
                            if "Rectangle" in item_name or "矩形" in item_name:
                                if item_name != gray_frame.Name:
                                    width_cm = item.Width / 72 * 2.54
                                    if width_cm > 15:  # 大于15cm的矩形
                                        white_frame = item
                                        break

                        # 记录原始位置
                        gray_frame_top = gray_frame.Top

                        # 获取文本实际高度
                        text_height_pt = text_box.TextFrame.TextRange.BoundHeight

                        # 计算边距
                        margin_top_pt = MARGIN_CM / 2.54 * 72
                        margin_bottom_pt = MARGIN_BOTTOM_CM / 2.54 * 72

                        # 新的文本框高度 = 文本高度 + 上边距 + 下边距
                        new_text_height_pt = text_height_pt + margin_top_pt + margin_bottom_pt

                        # 设置文本框边距
                        text_box.TextFrame.MarginTop = margin_top_pt
                        text_box.TextFrame.MarginBottom = margin_bottom_pt

                        # 设置文本框高度
                        text_box.Height = new_text_height_pt

                        # 计算新的灰色框高度
                        # 灰色框顶部不变，底部与文本框底部对齐
                        text_box_bottom = text_box.Top + text_box.Height
                        new_gray_height = text_box_bottom - gray_frame_top

                        # 设置灰色框高度
                        gray_frame.Height = new_gray_height

                        # 调整白色框高度，使其底部与灰色框顶部对齐
                        if white_frame:
                            white_frame_top = white_frame.Top
                            new_white_height = gray_frame_top - white_frame_top
                            if new_white_height > 0:
                                white_frame.Height = new_white_height
                                print(f"白色框调整: height={new_white_height/72*2.54:.2f}cm")

                        print("补充说明区域调整:")
                        print(f"  文本高度: {text_height_pt:.1f}pt")
                        print(f"  文本框新高度: {new_text_height_pt:.1f}pt")
                        print(f"  灰色框新高度: {new_gray_height:.1f}pt ({new_gray_height/72*2.54:.2f}cm)")
                        break

        pres.Save()
        pres.Close()

        return True

    except Exception as e:
        print(f"调整补充说明高度失败: {e}")
        import traceback
        traceback.print_exc()
        return False


# 注意: adjust_knowledge_content_height 函数已被废弃
# 边距设置现在在 _fill_knowledge_content 中通过 XML 直接完成
# 高度调整在 _adjust_shape_height_via_com 中通过公式计算完成
# 如果需要COM API后处理，请使用 distribute_groups_vertically 函数


def verify_series_replacement(presentation: Presentation, expected_name: str, expected_level: int, include_cover: bool = True) -> tuple:
    """
    验证PPT中的系列名称是否正确替换

    Args:
        presentation: PPT演示文稿对象
        expected_name: 期望的系列名称
        expected_level: 期望的阶数
        include_cover: 是否包含封面页。当为False时，跳过布局0的检查，
                       因为封面布局只在include_cover=True时才会被修改

    Returns:
        (是否成功, 错误列表)
    """
    errors = []

    if len(presentation.slide_masters) < 1:
        return False, ["没有找到幻灯片母版"]

    master = presentation.slide_masters[0]
    layouts = master.slide_layouts

    def check_shapes(shapes, location):
        for shape in shapes:
            if hasattr(shape, 'text') and shape.text:
                text = shape.text
                # 检查是否还有旧的模板文字
                if '机械臂设计' in text and expected_name != '机械臂设计':
                    errors.append(f'{location} - {shape.name}: 仍包含\"机械臂设计\" -> \"{text[:50]}\"')
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    check_shapes(shape.shapes, location)
                except:
                    pass

    # 检查布局0（封面页布局）- 只有包含封面时才检查，因为布局只在include_cover=True时才被修改
    if include_cover and len(layouts) > 0:
        check_shapes(layouts[0].shapes, '布局0(封面)')

    # 检查布局1（基本信息页布局）
    if len(layouts) > 1:
        check_shapes(layouts[1].shapes, '布局1(基本信息)')

    return len(errors) == 0, errors
